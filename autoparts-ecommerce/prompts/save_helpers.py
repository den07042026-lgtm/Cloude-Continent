"""
Вспомогательные функции сохранения документов — Autoparts Analytics System v1.
Используется всеми промптами. Вызывается Claude автоматически.

Адаптировано из IFOAM Analytics save_helpers.py для автозапчастей.
"""

import os
import datetime
from pathlib import Path
from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

# ─── Цвета — стиль "автомобильный / технический" ─────────────────────────────
DARK_RGB  = RGBColor(0x1A, 0x1A, 0x2E)   # тёмно-синий
ACCT_RGB  = RGBColor(0xE6, 0x3A, 0x00)   # красно-оранжевый (технический акцент)
GREY_RGB  = RGBColor(0x88, 0x88, 0x88)
DARK_HEX  = "1A1A2E"
ACCT_HEX  = "E63A00"
LIGHT_HEX = "FFF3EE"
WHITE_HEX = "FFFFFF"

BASE_FOLDER = Path("C:/Users/Admin/Downloads")


def get_product_folder(brand: str, product_name: str, article: str) -> Path:
    """Создаёт и возвращает папку продукта в Downloads."""
    safe = lambda s: "".join(c for c in str(s) if c.isalnum() or c in "-_").strip()
    folder_name = f"{safe(brand)}_{safe(product_name)}_{safe(article)}"
    folder = BASE_FOLDER / folder_name
    folder.mkdir(parents=True, exist_ok=True)
    return folder


def _set_cell_bg(cell, hex_color: str):
    """Устанавливает цвет фона ячейки таблицы Word."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def _set_cell_border(cell, **kwargs):
    """Устанавливает границы ячейки таблицы Word."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = OxmlElement("w:tcBorders")
    for side in ("top", "left", "bottom", "right", "insideH", "insideV"):
        val = kwargs.get(side, {"sz": 4, "val": "single", "color": "CCCCCC"})
        border = OxmlElement(f"w:{side}")
        border.set(qn("w:sz"), str(val.get("sz", 4)))
        border.set(qn("w:val"), val.get("val", "single"))
        border.set(qn("w:color"), val.get("color", "CCCCCC"))
        tcBorders.append(border)
    tcPr.append(tcBorders)


def make_docx_card(folder: Path, filename: str, title: str,
                   sections: list) -> str:
    """
    Создаёт DOCX с таблицами в структуре карточки запчасти.

    sections — список словарей:
      {
        "heading": "КАРТОЧКА ЗАПЧАСТИ",
        "type": "table",
        "rows": [("Поле", "Значение"), ...]
      }
      или
      {
        "heading": "Аудитория и УТП",
        "type": "text",
        "text": "строка"
      }
    """
    doc = Document()

    for sec in doc.sections:
        sec.top_margin    = Cm(1.5)
        sec.bottom_margin = Cm(1.5)
        sec.left_margin   = Cm(2.0)
        sec.right_margin  = Cm(1.5)

    h = doc.add_heading(title, level=0)
    h.runs[0].font.color.rgb = DARK_RGB
    h.runs[0].font.size = Pt(16)

    p = doc.add_paragraph(f"Дата создания: {datetime.date.today().strftime('%d.%m.%Y')}")
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = GREY_RGB
    doc.add_paragraph("")

    for sec in sections:
        if sec.get("heading"):
            h2 = doc.add_heading(sec["heading"], level=1)
            for run in h2.runs:
                run.font.color.rgb = ACCT_RGB
                run.font.size = Pt(12)
            doc.add_paragraph("")

        if sec.get("type") == "table" and sec.get("rows"):
            rows = sec["rows"]
            tbl = doc.add_table(rows=len(rows), cols=2)
            tbl.style = "Table Grid"
            tbl.autofit = False
            tbl.columns[0].width = Cm(6)
            tbl.columns[1].width = Cm(11)

            for i, (field, value) in enumerate(rows):
                cell_f = tbl.cell(i, 0)
                cell_f.text = str(field)
                run_f = cell_f.paragraphs[0].runs[0]
                run_f.font.bold = True
                run_f.font.size = Pt(10)
                run_f.font.color.rgb = DARK_RGB
                _set_cell_bg(cell_f, "F0F0F4")

                cell_v = tbl.cell(i, 1)
                cell_v.text = str(value)
                run_v = cell_v.paragraphs[0].runs[0]
                run_v.font.size = Pt(10)
                if i % 2 == 0:
                    _set_cell_bg(cell_v, WHITE_HEX)
                else:
                    _set_cell_bg(cell_v, LIGHT_HEX)

            doc.add_paragraph("")

        elif sec.get("type") == "text" and sec.get("text"):
            p = doc.add_paragraph(sec["text"])
            if p.runs:
                p.runs[0].font.size = Pt(11)
            doc.add_paragraph("")

    path = folder / filename
    doc.save(str(path))
    print(f"Сохранено: {path}")
    return str(path)


def make_docx(folder: Path, filename: str, title: str, content: str) -> str:
    """Создаёт DOCX с текстовым содержимым (для промптов без таблиц)."""
    doc = Document()
    for sec in doc.sections:
        sec.top_margin    = Cm(1.5)
        sec.bottom_margin = Cm(1.5)
        sec.left_margin   = Cm(2.0)
        sec.right_margin  = Cm(1.5)

    h = doc.add_heading(title, level=0)
    h.runs[0].font.color.rgb = DARK_RGB

    p = doc.add_paragraph(f"Дата создания: {datetime.date.today().strftime('%d.%m.%Y')}")
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = GREY_RGB
    doc.add_paragraph("")

    for line in content.strip().split("\n"):
        line = line.rstrip()
        if line.startswith("===") and line.endswith("==="):
            sec_text = line.strip("= ").strip()
            h2 = doc.add_heading(sec_text, level=1)
            for run in h2.runs:
                run.font.color.rgb = ACCT_RGB
        elif line.startswith("---"):
            doc.add_paragraph("─" * 60)
        elif line == "":
            doc.add_paragraph("")
        else:
            p = doc.add_paragraph(line)
            if p.runs:
                p.runs[0].font.size = Pt(11)

    path = folder / filename
    doc.save(str(path))
    print(f"Сохранено: {path}")
    return str(path)


def _col_letter(n: int) -> str:
    """Возвращает буквы колонки Excel по номеру (1-based)."""
    result = ""
    while n > 0:
        n, rem = divmod(n - 1, 26)
        result = chr(65 + rem) + result
    return result


def make_xlsx_competitors(folder: Path, filename: str, df) -> str:
    """Сохраняет таблицу конкурентов из pandas DataFrame в XLSX."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Конкуренты"

    headers = list(df.columns)
    n_cols = max(len(headers), 1)
    last_col = _col_letter(n_cols)

    ws.merge_cells(f"A1:{last_col}1")
    tc = ws["A1"]
    tc.value = "Анализ конкурентов — Автозапчасти"
    tc.font = Font(bold=True, size=14, color=WHITE_HEX)
    tc.fill = PatternFill("solid", fgColor=DARK_HEX)
    tc.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 28

    ws.merge_cells(f"A2:{last_col}2")
    dc = ws["A2"]
    dc.value = f"Дата: {datetime.date.today().strftime('%d.%m.%Y')}"
    dc.font = Font(size=9, color="888888")
    dc.alignment = Alignment(horizontal="right")

    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=4, column=col_idx, value=header)
        cell.font = Font(bold=True, color=WHITE_HEX)
        cell.fill = PatternFill("solid", fgColor=ACCT_HEX)
        cell.alignment = Alignment(horizontal="center", wrap_text=True)
        cell.border = Border(
            bottom=Side(style="thin", color=DARK_HEX),
            right=Side(style="thin", color="CCCCCC")
        )

    for row_idx, (_, row) in enumerate(df.iterrows(), start=5):
        fill_color = WHITE_HEX if row_idx % 2 == 1 else LIGHT_HEX
        for col_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.fill = PatternFill("solid", fgColor=fill_color)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = Border(right=Side(style="thin", color="EEEEEE"))

    for col in ws.columns:
        cells = [c for c in col if hasattr(c, "column_letter") and c.row >= 4]
        if not cells:
            continue
        max_len = max((len(str(c.value or "")) for c in cells), default=8)
        ws.column_dimensions[cells[0].column_letter].width = min(max_len + 4, 45)

    ws.freeze_panes = "A5"
    path = folder / filename
    wb.save(str(path))
    print(f"Сохранено: {path}")
    return str(path)


def make_xlsx_pricing(folder: Path, filename: str, df) -> str:
    """Сохраняет таблицу юнит-экономики (сценарии цен) в XLSX."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Юнит-экономика"

    headers = list(df.columns)
    n_cols = max(len(headers), 1)
    last_col = _col_letter(n_cols)

    ws.merge_cells(f"A1:{last_col}1")
    tc = ws["A1"]
    tc.value = "Юнит-экономика — сценарии ценообразования"
    tc.font = Font(bold=True, size=14, color=WHITE_HEX)
    tc.fill = PatternFill("solid", fgColor=DARK_HEX)
    tc.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 28

    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=3, column=col_idx, value=header)
        cell.font = Font(bold=True, color=WHITE_HEX)
        cell.fill = PatternFill("solid", fgColor=ACCT_HEX)
        cell.alignment = Alignment(horizontal="center", wrap_text=True)

    for row_idx, (_, row) in enumerate(df.iterrows(), start=4):
        margin = row.get("Маржа %", 0)
        for col_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            if margin >= 20:
                cell.fill = PatternFill("solid", fgColor="E8F5E9")   # зелёный
            elif margin >= 15:
                cell.fill = PatternFill("solid", fgColor="FFF8E1")   # жёлтый
            else:
                cell.fill = PatternFill("solid", fgColor="FFEBEE")   # красный
            cell.alignment = Alignment(horizontal="center")

    for col in ws.columns:
        cells = [c for c in col if hasattr(c, "column_letter") and c.row >= 3]
        if not cells:
            continue
        max_len = max((len(str(c.value or "")) for c in cells), default=8)
        ws.column_dimensions[cells[0].column_letter].width = min(max_len + 4, 30)

    ws.freeze_panes = "A4"
    path = folder / filename
    wb.save(str(path))
    print(f"Сохранено: {path}")
    return str(path)


def make_xlsx_seo(folder: Path, filename: str, df) -> str:
    """Сохраняет SEO-таблицу поисковых запросов в XLSX."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "SEO запросы"

    headers = list(df.columns)
    n_cols = max(len(headers), 1)
    last_col = _col_letter(n_cols)

    ws.merge_cells(f"A1:{last_col}1")
    tc = ws["A1"]
    tc.value = "SEO — поисковые запросы"
    tc.font = Font(bold=True, size=14, color=WHITE_HEX)
    tc.fill = PatternFill("solid", fgColor=DARK_HEX)
    tc.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 28

    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=3, column=col_idx, value=header)
        cell.font = Font(bold=True, color=WHITE_HEX)
        cell.fill = PatternFill("solid", fgColor=ACCT_HEX)
        cell.alignment = Alignment(horizontal="center", wrap_text=True)

    for row_idx, (_, row) in enumerate(df.iterrows(), start=4):
        fill_color = WHITE_HEX if row_idx % 2 == 0 else LIGHT_HEX
        for col_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.fill = PatternFill("solid", fgColor=fill_color)
            cell.alignment = Alignment(wrap_text=True, vertical="top")

    for col in ws.columns:
        cells = [c for c in col if hasattr(c, "column_letter") and c.row >= 3]
        if not cells:
            continue
        max_len = max((len(str(c.value or "")) for c in cells), default=8)
        ws.column_dimensions[cells[0].column_letter].width = min(max_len + 4, 60)

    ws.freeze_panes = "A4"
    path = folder / filename
    wb.save(str(path))
    print(f"Сохранено: {path}")
    return str(path)


def make_pptx_designer_tz(folder: Path, filename: str,
                           slides: list) -> str:
    """
    Создаёт PPTX-презентацию ТЗ для дизайнера.

    slides — список словарей:
      {"title": "Слайд 1 — Главное фото", "content": "описание содержания слайда"}
    """
    prs = Presentation()
    prs.slide_width  = Emu(9144000)   # 16:9
    prs.slide_height = Emu(5143500)

    blank_layout = prs.slide_layouts[6]  # пустой макет

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)

        # Фон слайда — тёмный
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PptRGB(0x1A, 0x1A, 0x2E)

        # Заголовок
        txBox = slide.shapes.add_textbox(
            PptInches(0.3), PptInches(0.2),
            PptInches(15.6), PptInches(0.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_data.get("title", "")
        run.font.size = PptPt(24)
        run.font.bold = True
        run.font.color.rgb = PptRGB(0xE6, 0x3A, 0x00)

        # Разделительная линия
        line = slide.shapes.add_connector(1, PptInches(0.3), PptInches(1.0),
                                           PptInches(15.6), PptInches(1.0))
        line.line.color.rgb = PptRGB(0xE6, 0x3A, 0x00)
        line.line.width = Emu(25400)

        # Содержимое
        txBox2 = slide.shapes.add_textbox(
            PptInches(0.3), PptInches(1.1),
            PptInches(15.6), PptInches(4.0)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        content = slide_data.get("content", "")
        for line_text in content.split("\n"):
            if tf2.paragraphs[0].text == "":
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            run2 = p2.add_run()
            run2.text = line_text
            run2.font.size = PptPt(14)
            run2.font.color.rgb = PptRGB(0xE8, 0xE8, 0xE8)

    path = folder / filename
    prs.save(str(path))
    print(f"Сохранено: {path}")
    return str(path)
