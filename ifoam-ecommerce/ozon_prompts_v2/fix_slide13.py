# -*- coding: utf-8 -*-
"""Заменяет слайд 13 (Экономика) — без рублей, только демонстрация экономии расхода."""
import sys, io, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

sys.path.insert(0, "C:/Users/1/Downloads/ozon_prompts_v2")
from save_helpers import get_product_folder
folder    = get_product_folder("ifoam", "OrangeStandard5kg", "119105")
pptx_path = folder / "07_designer_tz_119105_v2.pptx"
photos_dir = Path("C:/Users/1/Downloads/ifoam_OrangeStandard5kg_119105/Фото карточек конкурентов")
tmp_dir   = Path("C:/Users/1/Downloads/_tmp_png2")
tmp_dir.mkdir(exist_ok=True)

def to_png(stem):
    src = photos_dir / f"{stem}.webp"
    dst = tmp_dir / f"{stem}.png"
    if not dst.exists():
        img = Image.open(src).convert("RGB")
        img.save(dst, "PNG")
    return str(dst)

# ── цвета ─────────────────────────────────────────────────────
BLACK      = RGBColor(0x1A,0x1A,0x1A)
ORANGE     = RGBColor(0xFF,0x6B,0x00)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
LIGHT_GRAY = RGBColor(0xF2,0xF2,0xF2)
HDR_BG     = RGBColor(0x33,0x18,0x00)
GREEN      = RGBColor(0x2E,0xCC,0x71)
RED        = RGBColor(0xFF,0x44,0x44)
ORANGE_LITE= RGBColor(0xFF,0xA0,0x40)
DARK_GRAY  = RGBColor(0x2D,0x2D,0x2D)

prs   = Presentation(str(pptx_path))
BLANK = prs.slide_layouts[6]

# ── удалить старый слайд 13 (индекс 12) ────────────────────────
slide_idx = 12   # 0-based → слайд 13
xml_slides = prs.slides._sldIdLst
rId_to_remove = prs.slides._sldIdLst[slide_idx].get('r:id')

# Удаляем элемент из списка слайдов
sldIdLst = prs.slides._sldIdLst
elem = sldIdLst[slide_idx]
sldIdLst.remove(elem)

# ── добавить новый слайд в конец, затем переместить на позицию 12 ──
new_slide = prs.slides.add_slide(BLANK)

# Переместить в нужную позицию (последний → позиция 12)
sldIdLst = prs.slides._sldIdLst
# Извлекаем последний элемент
last_elem = sldIdLst[-1]
sldIdLst.remove(last_elem)
sldIdLst.insert(slide_idx, last_elem)

s = prs.slides[slide_idx]  # теперь это наш новый слайд

# ── фон ────────────────────────────────────────────────────────
fill = s.background.fill
fill.solid()
fill.fore_color.rgb = BLACK

def txbox(slide, text, left, top, width, height,
          font_size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txb

def rect(slide, left, top, width, height, fill_color, line_color=None, lw=0):
    from pptx.util import Pt as PPt
    sh = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = PPt(lw)
    else:
        sh.line.fill.background()
    return sh

def divider(slide, top, color=ORANGE):
    rect(slide, 0.2, top, 9.6, 0.04, fill_color=color)

def add_photo(slide, png_path, left, top, width, height):
    img = Image.open(png_path)
    iw, ih = img.size
    ratio = iw / ih
    box_ratio = width / height
    if ratio > box_ratio:
        w, h = width, width / ratio
        l, t = left, top + (height - h) / 2
    else:
        h, w = height, height * ratio
        l, t = left + (width - w) / 2, top
    slide.shapes.add_picture(png_path, Inches(l), Inches(t), Inches(w), Inches(h))

# ════════════════════════════════════════════════════════════════
# ЗАГОЛОВОК
# ════════════════════════════════════════════════════════════════
rect(s, 0.2, 0.12, 9.6, 0.72, fill_color=HDR_BG)
txbox(s, "ФОТО 4: ЭКОНОМИКА РАСХОДА — как конкуренты показывают концентрацию",
      0.3, 0.17, 9.4, 0.65,
      font_size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
divider(s, 0.88)

# ════════════════════════════════════════════════════════════════
# ФОТО КОНКУРЕНТОВ (4 штуки, 2 колонки × 2 ряда)
# ════════════════════════════════════════════════════════════════
photos = [
    ("8278219912",
     "DR.BERG ACE\nТаблица разбавления — пеногенератор\n1:100–200, пенокомплект 1:10–20",
     "DR.BERG"),
    ("8278220387",
     "DR.BERG ACE\n«Высокая концентрация средства»\n+ «эффективен при жёсткой воде»",
     "DR.BERG"),
    ("8901708053",
     "BIGHIM\n«5л = до 150 моек» крупно\nна главном фото — цифра выходит вперёд",
     "BIGHIM"),
    ("8278266808",
     "DR.BERG ACE\n«Автошампунь до 200 моек»\nна главном продакт-шоте",
     "DR.BERG"),
]

brand_colors = {
    "GLITTER": RGBColor(0x00,0xAA,0xDD),
    "DR.BERG": RGBColor(0x66,0xCC,0x22),
    "GRASS":   RGBColor(0x22,0xAA,0x33),
    "BIGHIM":  RGBColor(0xFF,0x88,0x00),
}

PHOTO_TOP    = 1.0
PHOTO_HEIGHT = 3.5
CAPTION_H    = 0.82
MARGIN, GAP  = 0.22, 0.18
cols = 2
photo_w = (10.0 - 2*MARGIN - GAP*(cols-1)) / cols   # ~4.7

for idx, (stem, caption, brand) in enumerate(photos):
    col = idx % cols
    row = idx // cols
    left = MARGIN + col * (photo_w + GAP)
    top  = PHOTO_TOP + row * (PHOTO_HEIGHT + CAPTION_H + 0.1)

    png = to_png(stem)
    rect(s, left, top, photo_w, PHOTO_HEIGHT + CAPTION_H + 0.08,
         fill_color=RGBColor(0x22,0x22,0x22),
         line_color=ORANGE, lw=1.0)
    add_photo(s, png, left+0.05, top+0.05, photo_w-0.1, PHOTO_HEIGHT-0.1)

    bc = brand_colors.get(brand, ORANGE)
    rect(s, left + photo_w - 1.3, top + 0.05, 1.22, 0.32, fill_color=bc)
    txbox(s, brand, left + photo_w - 1.3, top + 0.05, 1.22, 0.32,
          font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txbox(s, caption, left+0.06, top+PHOTO_HEIGHT+0.02,
          photo_w-0.12, CAPTION_H, font_size=11, color=LIGHT_GRAY)

# ════════════════════════════════════════════════════════════════
# НИЖНИЙ БЛОК — ФОКУС ДЛЯ ifoam
# ════════════════════════════════════════════════════════════════
focus_top = PHOTO_TOP + 2*(PHOTO_HEIGHT + CAPTION_H + 0.1) + 0.1
divider(s, focus_top)

txbox(s, "▶  НА ЧТО ОБРАТИТЬ ВНИМАНИЕ — ФОТО 4 ifoam Orange Standard",
      0.2, focus_top+0.1, 9.6, 0.45,
      font_size=15, bold=True, color=ORANGE)

# Три блока-иконки: 5 кг → 150–200 л → 20 авто
blocks = [
    ("5 кг\nконцентрата", "→"),
    ("150–200 л\nрабочего раствора", "→"),
    ("≥ 20\nавтомобилей", ""),
]
bw = 2.8
for i, (label, arrow) in enumerate(blocks):
    bx = 0.25 + i * (bw + 0.55)
    rect(s, bx, focus_top+0.65, bw, 1.1,
         fill_color=RGBColor(0x2A,0x14,0x00),
         line_color=ORANGE, lw=2)
    txbox(s, label, bx, focus_top+0.7, bw, 1.0,
          font_size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    if arrow:
        txbox(s, "→", bx+bw+0.05, focus_top+0.85, 0.45, 0.6,
              font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

focus_items = [
    ("✓", "Заголовок: «5 КГ = 20 АВТОМОБИЛЕЙ» — цифра выходит на первый план"),
    ("✓", "Визуальная цепочка: канистра → ведро с раствором → ряд автомобилей"),
    ("✓", "Таблица пропорций: пеногенератор 1:30–1:40 / пенокомплект 1:3–1:4"),
    ("✓", "Подпись: «1 кг концентрата = 30–40 литров рабочего состава»"),
    ("✗", "Без цен и рублей — только объём, расход, количество моек"),
]
for j, (mark, text) in enumerate(focus_items):
    col_mark = {" ✓": GREEN, "✓": GREEN, "✗": RED}.get(mark, GREEN)
    y = focus_top + 1.95 + j * 0.4
    txbox(s, mark, 0.2, y, 0.4, 0.4,
          font_size=14, bold=True, color=GREEN if mark=="✓" else RED,
          align=PP_ALIGN.CENTER)
    txbox(s, text, 0.65, y, 9.1, 0.4, font_size=13, color=WHITE)

txbox(s, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.2, 12.9, 9.6, 0.35, font_size=11,
      color=RGBColor(0x66,0x66,0x66), align=PP_ALIGN.CENTER)

# ── сохранить ─────────────────────────────────────────────────
prs.save(str(pptx_path))
print(f"Сохранено: {pptx_path}  |  Слайдов: {len(prs.slides)}")
shutil.rmtree(tmp_dir, ignore_errors=True)
