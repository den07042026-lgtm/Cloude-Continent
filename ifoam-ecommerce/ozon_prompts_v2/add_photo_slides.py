# -*- coding: utf-8 -*-
"""
Добавляет слайды с реальными фото конкурентов в ТЗ для дизайнера.
По одному слайду на каждый тип нашего фото: 2-4 фото конкурента + аннотации.
"""
import sys, io, os, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── пути ──────────────────────────────────────────────────────
sys.path.insert(0, "C:/Users/1/Downloads/ozon_prompts_v2")
from save_helpers import get_product_folder
folder    = get_product_folder("ifoam", "OrangeStandard5kg", "119105")
pptx_path = folder / "07_designer_tz_119105_v2.pptx"
photos_dir = Path("C:/Users/1/Downloads/ifoam_OrangeStandard5kg_119105/Фото карточек конкурентов")
tmp_dir    = Path("C:/Users/1/Downloads/_tmp_png")
tmp_dir.mkdir(exist_ok=True)

# ── конвертация webp → png ─────────────────────────────────────
def to_png(stem):
    src = photos_dir / f"{stem}.webp"
    dst = tmp_dir / f"{stem}.png"
    if not dst.exists():
        img = Image.open(src).convert("RGB")
        img.save(dst, "PNG")
    return str(dst)

# ── цвета ──────────────────────────────────────────────────────
BLACK      = RGBColor(0x1A,0x1A,0x1A)
ORANGE     = RGBColor(0xFF,0x6B,0x00)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
LIGHT_GRAY = RGBColor(0xF2,0xF2,0xF2)
DARK_GRAY  = RGBColor(0x2D,0x2D,0x2D)
ORANGE_LITE= RGBColor(0xFF,0xA0,0x40)
HDR_BG     = RGBColor(0x33,0x18,0x00)
GREEN      = RGBColor(0x2E,0xCC,0x71)
YELLOW     = RGBColor(0xFF,0xCC,0x00)
RED        = RGBColor(0xFF,0x44,0x44)

prs  = Presentation(str(pptx_path))
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          font_size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Arial"
    return txb

def rect(slide, left, top, width, height, fill_color, line_color=None, lw=0):
    from pptx.util import Pt as PPt
    s = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb  = line_color
        s.line.width = PPt(lw)
    else:
        s.line.fill.background()
    return s

def divider(slide, top, color=ORANGE):
    rect(slide, 0.2, top, 9.6, 0.04, fill_color=color)

def add_photo(slide, png_path, left, top, width, height):
    """Вставить фото с сохранением пропорций внутри прямоугольника."""
    from pptx.util import Inches
    img = Image.open(png_path)
    iw, ih = img.size
    ratio = iw / ih
    box_ratio = width / height
    if ratio > box_ratio:
        w = width
        h = width / ratio
        t = top + (height - h) / 2
        l = left
    else:
        h = height
        w = height * ratio
        l = left + (width - w) / 2
        t = top
    slide.shapes.add_picture(png_path, Inches(l), Inches(t), Inches(w), Inches(h))

# ════════════════════════════════════════════════════════════════
# Данные: (наш_слайд, заголовок, фокус_буллеты, [(stem, подпись, бренд), ...])
# ════════════════════════════════════════════════════════════════
slides_data = [

    # ── 1. ГЛАВНОЕ ФОТО ───────────────────────────────────────
    {
        "num": 1, "title": 'ФОТО 1: ГЛАВНОЕ ФОТО — примеры конкурентов',
        "focus": [
            ("✓", "Тёмный фон (#1A1A1A) — отстройка от светлого GLITTER G10"),
            ("✓", "Оранжевый акцент — цвет жидкости виден через канистру"),
            ("✓", "Этикетка Orange Standard читается полностью"),
            ("✓", "Бейдж «5 кг» + метка «ОТ ПРОИЗВОДИТЕЛЯ»"),
            ("✗", "Не белый фон — на полке сливаемся с GLITTER G10"),
        ],
        "photos": [
            ("6592383288",  "GLITTER G10\nЧистый продакт-шот,\nсветлый фон, этикетка читается",      "GLITTER"),
            ("8278266808",  "DR.BERG ACE\nТёмный фон, канистр крупно,\n«до 200 моек» — цифры видны", "DR.BERG"),
            ("9463559615",  "GRASS Balance\n«Бренд №1» бейдж, яркий\nжёлтый — читается сразу",       "GRASS"),
            ("8901708053",  "BIGHIM\n«5л = до 150 моек» крупно,\nяркий оранжевый фон",               "BIGHIM"),
        ],
    },

    # ── 2. ХАРАКТЕРИСТИКИ / УТП ───────────────────────────────
    {
        "num": 2, "title": 'ФОТО 2: ХАРАКТЕРИСТИКИ / УТП — примеры конкурентов',
        "focus": [
            ("✓", "pH 12 — крупный блок отдельно (никто из конкурентов так НЕ делает)"),
            ("✓", "Пропорции 1:30–1:40 / 1:3–1:4 — иконки или таблица"),
            ("✓", "«Не повреждает ЛКП» — закрываем главное возражение"),
            ("✗", "Не зелёные маркеры — цвет занят брендом GRASS"),
            ("✗", "Не более 5–6 буллетов — не перегружать"),
        ],
        "photos": [
            ("7388654580",  "GLITTER G10\n4 буллета с иконками:\nпена / вода / смывается / разводы",  "GLITTER"),
            ("8278219581",  "DR.BERG ACE\n3 плашки на тёмном:\nактивная пена / плотная / безопасно",  "DR.BERG"),
            ("9405795810",  "GLITTER G5\n«Готов к применению»\nбуллеты + логотип",                   "GLITTER"),
            ("8901704980",  "BIGHIM\nСписок на ярком фоне:\nкислотно-щелочной / блеск / водоотталк.", "BIGHIM"),
        ],
    },

    # ── 3. ИНСТРУКЦИЯ ─────────────────────────────────────────
    {
        "num": 3, "title": 'ФОТО 3: ИНСТРУКЦИЯ ПО ПРИМЕНЕНИЮ — примеры конкурентов',
        "focus": [
            ("✓", "4 шага: остудите / разведите / нанесите / смойте"),
            ("✓", "Шаг 2: КРУПНО пропорции — 1:30–1:40 и 1:3–1:4"),
            ("✓", "Реальные фото процесса (не иконки) — как у GRASS"),
            ("✓", "Предупреждение «НЕ ДОПУСКАЙТЕ ВЫСЫХАНИЯ» жёлтым"),
            ("✗", "Не мелкий текст — пропорции должны читаться с телефона"),
        ],
        "photos": [
            ("7448143493",  "GLITTER G10\nИнструкция с иконками:\nпеногенератор / пенокомплект / дозатрон", "GLITTER"),
            ("9405795764",  "GLITTER G5\n«4 шага до чистого кузова»\nСетка 2×2 с реальными фото",          "GLITTER"),
            ("9463559637",  "GRASS Balance\nШаг 1: «Откройте упаковку»\nРеальные руки, тёмный стиль",      "GRASS"),
            ("9463559647",  "GRASS Balance\nШаг 2: «Разведите состав»\nПоказан концентрат наглядно",       "GRASS"),
        ],
    },

    # ── 4. ЭКОНОМИКА ──────────────────────────────────────────
    {
        "num": 4, "title": 'ФОТО 4: ЭКОНОМИКА — как конкуренты показывают выгоду',
        "focus": [
            ("✓", "УНИКАЛЬНЫЙ слайд — у конкурентов нет отдельной «экономики»"),
            ("✓", "Заголовок: «5 КГ = 20 АВТОМОБИЛЕЙ» или «МОЙКА ОТ 12 РУБЛЕЙ»"),
            ("✓", "Три цифровых блока: 5 кг → 150–200 л → ~12 руб/мойка"),
            ("✓", "Таблица сравнения: автомойка 300–600р vs ifoam ~12р"),
            ("✗", "Не перегружать расчётами — только 3 ключевые цифры крупно"),
        ],
        "photos": [
            ("8278219912",  "DR.BERG ACE\n«Разбавление»: ратио 1:100–200\nпеногенератор / пенокомплект",   "DR.BERG"),
            ("8278220387",  "DR.BERG ACE\n«Высокая концентрация»\n— ближайший аналог экономики",           "DR.BERG"),
            ("8278266808",  "DR.BERG ACE\n«До 200 моек» — мелко\nна главном фото",                        "DR.BERG"),
            ("8901708053",  "BIGHIM\n«5л = до 150 моек»\n— единственный кто считает",                     "BIGHIM"),
        ],
    },

    # ── 5. РЕКОМЕНДАЦИИ ───────────────────────────────────────
    {
        "num": 5, "title": 'ФОТО 5: РЕКОМЕНДАЦИИ / БЕЗОПАСНОСТЬ — примеры конкурентов',
        "focus": [
            ("✓", "Реальное фото человека в перчатках (не иконка) — как у GLITTER G5"),
            ("✓", "Главное предупреждение КРУПНО: «НЕ НАНОСИТЬ НА ГОРЯЧИЙ КУЗОВ»"),
            ("✓", "«ВЫДЕРЖАТЬ НЕ БОЛЕЕ 2 МИНУТ» — жёлтый акцент"),
            ("✓", "Блок антивозражений: «что будет если...» — закрываем страхи"),
            ("✗", "Не только список ⚠ — нужна визуальная иерархия по важности"),
        ],
        "photos": [
            ("6664644908",  "GLITTER G10/G5\n«Наденьте перчатки» крупно\nФото синей перчатки",        "GLITTER"),
            ("9405795820",  "GLITTER G5\n«Защитные перчатки»\nМужчина на мойке с пистолетом",         "GLITTER"),
            ("8278220086",  "DR.BERG ACE\n«Способ применения» + предупреждения\nТёмный, чёткий текст", "DR.BERG"),
        ],
    },

    # ── 6. ПРОИЗВОДСТВО / ДОВЕРИЕ ────────────────────────────
    {
        "num": 6, "title": 'ФОТО 6: ПРОИЗВОДСТВО И ДОВЕРИЕ — примеры конкурентов',
        "focus": [
            ("✓", "«ПРОИЗВОДИМ САМИ — ОТВЕЧАЕМ ЗА КАЧЕСТВО» — фото цеха ifoam"),
            ("✓", "Флаг РФ + ТУ 20.41.32-003-36603872-2024 — закрываем «подделка?»"),
            ("✓", "Коллаж 2×2: оборудование / лаборатория / розлив / готовый продукт"),
            ("✓", "Если есть отзывы — 2 карточки с цитатами и ★★★★★"),
            ("✗", "НЕ AI-фото лаборатории — только реальные кадры производства ifoam"),
        ],
        "photos": [
            ("9405795888",  "GLITTER G5\n«Сделано в России для наших дорог»\nКоллаж лаборатории 2×2",   "GLITTER"),
            ("9405795838",  "GLITTER G5\n«51 000+ положительных отзывов»\nКарточки с цитатами и ★",    "GLITTER"),
            ("9405795878",  "GLITTER G5\n«Мойка своими руками — как у профи»\nМужчина с канистром",    "GLITTER"),
            ("9405795857",  "GLITTER G5\n«Почему мы?» — таблица сравнения\nGLITTER vs обычное средство","GLITTER"),
        ],
    },
]

# ════════════════════════════════════════════════════════════════
# Генерация слайдов
# ════════════════════════════════════════════════════════════════
for data in slides_data:
    photos = data["photos"]
    n_photos = len(photos)

    # Сколько фото помещается в ряд: 3 или 4
    cols = min(n_photos, 4)
    rows_count = (n_photos + cols - 1) // cols  # 1 или 2

    # ── создаём слайд ───────────────────────────────────────
    s = prs.slides.add_slide(BLANK)
    bg(s, BLACK)

    # Заголовок
    rect(s, 0.2, 0.12, 9.6, 0.72, fill_color=HDR_BG)
    txbox(s, data["title"], 0.3, 0.17, 9.4, 0.65,
          font_size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    divider(s, 0.88)

    # ── фото конкурентов ────────────────────────────────────
    PHOTO_TOP    = 1.0          # верхний край фото
    PHOTO_HEIGHT = 3.8          # высота фото
    CAPTION_H    = 0.85         # высота подписи
    MARGIN       = 0.22         # отступ от краёв слайда
    GAP          = 0.15         # зазор между фото

    total_w = 10.0 - 2 * MARGIN
    photo_w = (total_w - GAP * (cols - 1)) / cols

    for idx, (stem, caption, brand) in enumerate(photos):
        col = idx % cols
        row = idx // cols

        left = MARGIN + col * (photo_w + GAP)
        top  = PHOTO_TOP + row * (PHOTO_HEIGHT + CAPTION_H + 0.1)

        png = to_png(stem)

        # Рамка-плашка
        rect(s, left, top, photo_w, PHOTO_HEIGHT + CAPTION_H + 0.08,
             fill_color=RGBColor(0x22,0x22,0x22),
             line_color=ORANGE, lw=1.0)

        # Фото
        add_photo(s, png, left + 0.05, top + 0.05, photo_w - 0.1, PHOTO_HEIGHT - 0.1)

        # Бренд-метка
        brand_colors = {
            "GLITTER": RGBColor(0x00,0xAA,0xDD),
            "DR.BERG": RGBColor(0x66,0xCC,0x22),
            "GRASS":   RGBColor(0x22,0xAA,0x33),
            "BIGHIM":  RGBColor(0xFF,0x88,0x00),
        }
        brand_col = brand_colors.get(brand, ORANGE)
        rect(s, left + photo_w - 1.3, top + 0.05, 1.22, 0.32,
             fill_color=brand_col)
        txbox(s, brand, left + photo_w - 1.3, top + 0.05, 1.22, 0.32,
              font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Подпись
        txbox(s, caption,
              left + 0.06, top + PHOTO_HEIGHT + 0.02,
              photo_w - 0.12, CAPTION_H,
              font_size=11, color=LIGHT_GRAY)

    # ── фокус для ifoam — снизу ─────────────────────────────
    # Определяем y начала блока фокуса
    used_rows = (n_photos + cols - 1) // cols
    focus_top = PHOTO_TOP + used_rows * (PHOTO_HEIGHT + CAPTION_H + 0.1) + 0.15

    divider(s, focus_top)

    txbox(s, f"▶  НА ЧТО ОБРАТИТЬ ВНИМАНИЕ — ФОТО {data['num']} ifoam Orange Standard",
          0.2, focus_top + 0.1, 9.6, 0.45,
          font_size=15, bold=True, color=ORANGE)

    focus_colors = {"✓": GREEN, "✗": RED, "→": ORANGE_LITE}
    available_h  = 13.33 - focus_top - 0.65
    line_h       = min(0.42, available_h / max(len(data["focus"]), 1))

    for j, (mark, text) in enumerate(data["focus"]):
        col_mark = focus_colors.get(mark, WHITE)
        y = focus_top + 0.6 + j * line_h
        txbox(s, mark, 0.2, y, 0.4, line_h,
              font_size=14, bold=True, color=col_mark, align=PP_ALIGN.CENTER)
        txbox(s, text, 0.65, y, 9.1, line_h,
              font_size=13, color=WHITE)

    # Нижний колонтитул
    txbox(s, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
          0.2, 12.9, 9.6, 0.35, font_size=11,
          color=RGBColor(0x66,0x66,0x66), align=PP_ALIGN.CENTER)

# ── сохранить ────────────────────────────────────────────────
prs.save(str(pptx_path))
print(f"Сохранено: {pptx_path}")
print(f"Итого слайдов: {len(prs.slides)}")

# Удалить временные PNG
shutil.rmtree(tmp_dir, ignore_errors=True)
print("Временные файлы удалены.")
