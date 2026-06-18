"""
Вспомогательные функции для сохранения документов по шагам анализа.
Используется всеми промптами. Вызывается Claude автоматически.
"""

import os
import datetime
from pathlib import Path
from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

# ─── Цвета бренда ────────────────────────────────────────────────────────────
DARK_RGB   = RGBColor(0x1A, 0x1A, 0x1A)
GOLD_RGB   = RGBColor(0xF4, 0xC0, 0x30)
GREY_RGB   = RGBColor(0x88, 0x88, 0x88)
DARK_HEX   = "1A1A1A"
GOLD_HEX   = "F4C030"
LIGHT_HEX  = "FFF8E1"
WHITE_HEX  = "FFFFFF"


def get_product_folder(brand: str, product_name: str, article: str) -> Path:
    """Создаёт и возвращает папку продукта."""
    safe = lambda s: "".join(c for c in str(s) if c.isalnum() or c in "-_").strip()
    folder_name = f"{safe(brand)}_{safe(product_name)}_{safe(article)}"
    folder = Path("C:/Users/1/Downloads") / folder_name
    folder.mkdir(parents=True, exist_ok=True)
    return folder


def _set_cell_bg(cell, hex_color: str):
    """Устанавливает цвет фона ячейки таблицы Word."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def _set_cell_border(cell, **kwargs):
    """Устанавливает границы ячейки таблицы Word."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = OxmlElement("w:tcBorders")
    for side in ("top", "left", "bottom", "right", "insideH", "insideV"):
        val = kwargs.get(side, {"sz": 4, "val": "single", "color": "CCCCCC"})
        border = OxmlElement(f"w:{side}")
        border.set(qn("w:sz"), str(val.get("sz", 4)))
        border.set(qn("w:val"), val.get("val", "single"))
        border.set(qn("w:color"), val.get("color", "CCCCCC"))
        tcBorders.append(border)
    tcPr.append(tcBorders)


def make_docx_card(folder: Path, filename: str, title: str,
                   sections: list) -> str:
    """
    Создаёт DOCX с таблицами в структуре карточки товара.

    sections — список словарей:
      {
        "heading": "КАРТОЧКА ТОВАРА",   # заголовок секции (опционально)
        "type": "table",                # "table" или "text"
        "rows": [("Поле", "Значение"), ...]  # для type="table"
        "text": "строка текста"         # для type="text"
      }
    """
    doc = Document()

    # Поля страницы
    for section in doc.sections:
        section.top_margin    = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        section.left_margin   = Cm(2.0)
        section.right_margin  = Cm(1.5)

    # Заголовок документа
    h = doc.add_heading(title, level=0)
    h.runs[0].font.color.rgb = DARK_RGB
    h.runs[0].font.size = Pt(16)

    # Дата
    p = doc.add_paragraph(f"Дата создания: {datetime.date.today().strftime('%d.%m.%Y')}")
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = GREY_RGB
    doc.add_paragraph("")

    for sec in sections:
        # Заголовок секции
        if sec.get("heading"):
            h2 = doc.add_heading(sec["heading"], level=1)
            for run in h2.runs:
                run.font.color.rgb = GOLD_RGB
                run.font.size = Pt(12)
            doc.add_paragraph("")

        if sec.get("type") == "table" and sec.get("rows"):
            rows = sec["rows"]
            tbl = doc.add_table(rows=len(rows), cols=2)
            tbl.style = "Table Grid"
            tbl.autofit = False
            tbl.columns[0].width = Cm(6)
            tbl.columns[1].width = Cm(11)

            for i, (field, value) in enumerate(rows):
                # Ячейка поля (левая)
                cell_f = tbl.cell(i, 0)
                cell_f.text = str(field)
                run_f = cell_f.paragraphs[0].runs[0]
                run_f.font.bold = True
                run_f.font.size = Pt(10)
                run_f.font.color.rgb = DARK_RGB
                _set_cell_bg(cell_f, "F5F5F5")

                # Ячейка значения (правая)
                cell_v = tbl.cell(i, 1)
                cell_v.text = str(value)
                run_v = cell_v.paragraphs[0].runs[0]
                run_v.font.size = Pt(10)
                if i % 2 == 0:
                    _set_cell_bg(cell_v, WHITE_HEX)
                else:
                    _set_cell_bg(cell_v, LIGHT_HEX)

            doc.add_paragraph("")

        elif sec.get("type") == "text" and sec.get("text"):
            p = doc.add_paragraph(sec["text"])
            p.runs[0].font.size = Pt(11)
            doc.add_paragraph("")

    path = folder / filename
    doc.save(str(path))
    return str(path)


def make_docx(folder: Path, filename: str, title: str, content: str) -> str:
    """Создаёт DOCX с текстовым содержимым (для промптов без таблиц)."""
    doc = Document()
    for section in doc.sections:
        section.top_margin    = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        section.left_margin   = Cm(2.0)
        section.right_margin  = Cm(1.5)

    h = doc.add_heading(title, level=0)
    h.runs[0].font.color.rgb = DARK_RGB

    p = doc.add_paragraph(f"Дата создания: {datetime.date.today().strftime('%d.%m.%Y')}")
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = GREY_RGB
    doc.add_paragraph("")

    for line in content.strip().split("\n"):
        line = line.rstrip()
        if line.startswith("===") and line.endswith("==="):
            sec_text = line.strip("= ").strip()
            h2 = doc.add_heading(sec_text, level=1)
            for run in h2.runs:
                run.font.color.rgb = GOLD_RGB
        elif line.startswith("---"):
            doc.add_paragraph("─" * 60)
        elif line == "":
            doc.add_paragraph("")
        else:
            p = doc.add_paragraph(line)
            if p.runs:
                p.runs[0].font.size = Pt(11)

    path = folder / filename
    doc.save(str(path))
    return str(path)


def _col_letter(n: int) -> str:
    """Возвращает букву(ы) колонки Excel по номеру (1-based). Поддерживает >26 колонок."""
    result = ""
    while n > 0:
        n, rem = divmod(n - 1, 26)
        result = chr(65 + rem) + result
    return result


def make_xlsx_table(folder: Path, filename: str, title: str,
                    headers: list, rows: list) -> str:
    """Создаёт XLSX с таблицей данных."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:31]

    last_col = _col_letter(max(len(headers), 1))

    # Заголовок листа
    ws.merge_cells(f"A1:{last_col}1")
    tc = ws["A1"]
    tc.value = title
    tc.font = Font(bold=True, size=14, color=WHITE_HEX)
    tc.fill = PatternFill("solid", fgColor=DARK_HEX)
    tc.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 28

    ws.merge_cells(f"A2:{last_col}2")
    dc = ws["A2"]
    dc.value = f"Дата: {datetime.date.today().strftime('%d.%m.%Y')}"
    dc.font = Font(size=9, color="888888")
    dc.alignment = Alignment(horizontal="right")

    # Шапка
    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=4, column=col_idx, value=header)
        cell.font = Font(bold=True, color=WHITE_HEX)
        cell.fill = PatternFill("solid", fgColor=GOLD_HEX)
        cell.alignment = Alignment(horizontal="center", wrap_text=True)
        cell.border = Border(
            bottom=Side(style="thin", color=DARK_HEX),
            right=Side(style="thin", color="CCCCCC")
        )

    # Данные
    for row_idx, row in enumerate(rows, start=5):
        fill_color = WHITE_HEX if row_idx % 2 == 1 else LIGHT_HEX
        for col_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.fill = PatternFill("solid", fgColor=fill_color)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = Border(right=Side(style="thin", color="EEEEEE"))

    # Авторазмер колонок
    for col in ws.columns:
        cells = [c for c in col if hasattr(c, "column_letter") and c.row >= 4]
        if not cells:
            continue
        max_len = max((len(str(c.value or "")) for c in cells), default=8)
        ws.column_dimensions[cells[0].column_letter].width = min(max_len + 4, 45)

    ws.freeze_panes = "A5"
    path = folder / filename
    wb.save(str(path))
    return str(path)


def make_xlsx_competitors(folder: Path, filename: str, df) -> str:
    """Сохраняет таблицу конкурентов из pandas DataFrame в XLSX."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Конкуренты"

    headers = list(df.columns)
    n_cols = max(len(headers), 1)
    last_col = _col_letter(n_cols)

    # Заголовок листа
    ws.merge_cells(f"A1:{last_col}1")
    tc = ws["A1"]
    tc.value = "Анализ конкурентов"
    tc.font = Font(bold=True, size=14, color=WHITE_HEX)
    tc.fill = PatternFill("solid", fgColor=DARK_HEX)
    tc.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 28

    ws.merge_cells(f"A2:{last_col}2")
    dc = ws["A2"]
    dc.value = f"Дата: {datetime.date.today().strftime('%d.%m.%Y')}"
    dc.font = Font(size=9, color="888888")
    dc.alignment = Alignment(horizontal="right")

    # Шапка
    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=4, column=col_idx, value=header)
        cell.font = Font(bold=True, color=WHITE_HEX)
        cell.fill = PatternFill("solid", fgColor=GOLD_HEX)
        cell.alignment = Alignment(horizontal="center", wrap_text=True)
        cell.border = Border(
            bottom=Side(style="thin", color=DARK_HEX),
            right=Side(style="thin", color="CCCCCC")
        )

    # Данные
    for row_idx, (_, row) in enumerate(df.iterrows(), start=5):
        fill_color = WHITE_HEX if row_idx % 2 == 1 else LIGHT_HEX
        for col_idx, value in enumerate(row, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.fill = PatternFill("solid", fgColor=fill_color)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = Border(right=Side(style="thin", color="EEEEEE"))

    # Авторазмер колонок
    for col in ws.columns:
        cells = [c for c in col if hasattr(c, "column_letter") and c.row >= 4]
        if not cells:
            continue
        max_len = max((len(str(c.value or "")) for c in cells), default=8)
        ws.column_dimensions[cells[0].column_letter].width = min(max_len + 4, 45)

    ws.freeze_panes = "A5"
    path = folder / filename
    wb.save(str(path))
    return str(path)


def make_pptx_designer_tz(folder: Path, filename: str,
                           product_info: dict, global_reqs: dict,
                           cards: list) -> str:
    """
    Создаёт PPTX-презентацию ТЗ для дизайнера.

    product_info: {"brand": str, "name": str, "article": str}
    global_reqs:  {"bg_color": "#0D1F18", "accent": "#F0C030", "text_color": "#FFFFFF", ...}
    cards: [{"number": int, "title": str, "priority": str,
             "task": str, "image": str, "text_on_photo": str, "reference": str}, ...]
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def _hex_to_ppt(hex_str: str) -> PptRGB:
        h = hex_str.strip("#")
        return PptRGB(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    bg_hex  = global_reqs.get("bg_color", "#1A1A1A")
    acc_hex = global_reqs.get("accent",   "#F4C030")
    BG   = _hex_to_ppt(bg_hex)
    GOLD = _hex_to_ppt(acc_hex)
    # Авто-выбор цвета текста: тёмный фон → белый, светлый → тёмный
    r, g, b = int(bg_hex.strip("#")[0:2], 16), int(bg_hex.strip("#")[2:4], 16), int(bg_hex.strip("#")[4:6], 16)
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    WHT = PptRGB(0xFF, 0xFF, 0xFF) if luminance < 0.5 else PptRGB(0x1A, 0x1A, 0x1A)
    GRY = PptRGB(0xAA, 0xAA, 0xAA) if luminance < 0.5 else PptRGB(0x66, 0x66, 0x66)

    def blank_slide():
        layout = prs.slide_layouts[6]   # пустой макет
        s = prs.slides.add_slide(layout)
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = BG
        return s

    def add_text(slide, text, left, top, width, height,
                 size=18, bold=False, color=WHT, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txb

    # ── Слайд 1: Технические требования ─────────────────────────────────────
    s1 = blank_slide()
    add_text(s1, "ТЕХНИЧЕСКИЕ ТРЕБОВАНИЯ", 0.4, 0.2, 8, 0.6, size=24, bold=True, color=GOLD)
    specs = [
        f"Форматы: 1200×1200 px (обложка) | 1200×1600 px (инфографика)",
        f"Качество: 300 dpi | RGB",
        f"Фон: {global_reqs.get('bg_color', '#1A1A1A')}  |  Акцент: {global_reqs.get('accent', '#F4C030')}",
        f"Типографика: жирный, заглавные, минимум 24px",
        f"Запрещено: лица людей, стоковые фото, текст поверх товара",
    ]
    for i, line in enumerate(specs):
        add_text(s1, line, 0.4, 1.0 + i * 0.65, 12, 0.6, size=16, color=WHT)

    # ── Слайд 2: Титульный ──────────────────────────────────────────────────
    s2 = blank_slide()
    brand   = product_info.get("brand", "")
    name    = product_info.get("name", "")
    article = product_info.get("article", "")
    add_text(s2, f"{brand} | {name}", 1, 2.2, 11, 1, size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s2, f"Арт. {article}", 1, 3.3, 11, 0.6, size=22, color=GRY, align=PP_ALIGN.CENTER)
    add_text(s2, "ТЗ ДЛЯ ДИЗАЙНЕРА", 1, 4.0, 11, 0.6, size=18, color=WHT, align=PP_ALIGN.CENTER)

    # ── Слайды карточек ─────────────────────────────────────────────────────
    PRIORITY_COLOR = {"срочно": PptRGB(0xFF, 0x55, 0x55),
                      "важно":  PptRGB(0xFF, 0xAA, 0x00),
                      "стандартно": GRY}

    for card in cards:
        s = blank_slide()
        num      = card.get("number", "?")
        title    = card.get("title", "")
        priority = card.get("priority", "стандартно").lower()
        pcolor   = PRIORITY_COLOR.get(priority, GRY)

        # Номер + название
        add_text(s, f"Карточка {num}", 0.3, 0.15, 3, 0.5, size=13, color=GRY)
        add_text(s, title.upper(), 0.3, 0.5, 9, 0.7, size=26, bold=True, color=GOLD)
        add_text(s, f"● {priority.upper()}", 9.8, 0.5, 3, 0.5, size=14, bold=True, color=pcolor)

        # Разделитель
        ln = s.shapes.add_connector(1, Inches(0.3), Inches(1.3), Inches(13), Inches(1.3))
        ln.line.color.rgb = GOLD
        ln.line.width = PptPt(1)

        # Левый блок — задача и изображение
        add_text(s, "ЗАДАЧА", 0.3, 1.45, 5.8, 0.4, size=11, bold=True, color=GOLD)
        add_text(s, card.get("task", ""), 0.3, 1.85, 5.8, 1.5, size=13, color=WHT)

        add_text(s, "ИЗОБРАЖЕНИЕ", 0.3, 3.45, 5.8, 0.4, size=11, bold=True, color=GOLD)
        add_text(s, card.get("image", ""), 0.3, 3.85, 5.8, 2.6, size=12, color=WHT)

        # Правый блок — текст и референс
        add_text(s, "ТЕКСТ НА ФОТО", 7.0, 1.45, 6, 0.4, size=11, bold=True, color=GOLD)
        add_text(s, card.get("text_on_photo", ""), 7.0, 1.85, 6, 2.0, size=12, color=WHT)

        add_text(s, "РЕФЕРЕНС", 7.0, 3.95, 6, 0.4, size=11, bold=True, color=GOLD)
        add_text(s, card.get("reference", ""), 7.0, 4.35, 6, 2.1, size=11, color=GRY)

    path = folder / filename
    prs.save(str(path))
    return str(path)


def make_pptx_designer_tz_v2(folder: Path, filename: str,
                              product_info: dict,
                              global_reqs: dict,
                              cards: list,
                              competitor_table: list = None,
                              photo_slides: list = None) -> str:
    """
    Создаёт PPTX-презентацию ТЗ для дизайнера по стандарту шаблона 119105.
    Формат: вертикальный 25.4×33.9 см.

    product_info: {"brand": str, "name": str, "article": str}
    global_reqs:  {
        "bg_color":    "#1A1A1A",
        "accent":      "#E91E8C",
        "text_color":  "#F2F2F2",
        "formats":     "1200×1600 px (инфографика) | 1200×1200 px (обложка)",
        "palette_note": "Пастельные фото: белый/розовый/лавандовый",
        "forbidden":   ["пункт1", "пункт2", ...],
        "utp_bullets": ["50 стирок из 1л", ...],
    }
    cards: [
      {
        "number":   int,
        "title":    "КАРТОЧКА ...",
        "goal":     "Цель карточки",
        "bg":       "Описание фона",
        "composition": "Описание композиции",
        "text_blocks": [("Метка", "Описание"), ...],
        "competitor_ref": "Справка по конкурентам",
        "do_dont": [("✓"/"✗", "текст"), ...],
      }, ...
    ]
    competitor_table: [
      {"num": 1, "photo": "Главное фото", "competitors": "...", "focus": "..."}, ...
    ]
    photo_slides: [
      {
        "photo_num": 1,
        "title": "ФОТО 1: ... — примеры конкурентов",
        "competitors": [("Бренд", "Описание"), ...],  # 4 шт
        "notes": [("✓"/"✗", "текст"), ...],
      }, ...
    ]

    Структура слайдов:
      Слайд 1:        Общие требования
      Слайды 2..N+1:  Карточки (N штук)
      Слайд N+2:      Технические требования
      Слайд N+3:      Анализ конкурентов (таблица)
      Слайды N+4..:   Фото 1..M — примеры конкурентов
    """
    from pptx import Presentation
    from pptx.util import Cm, Pt as PptPt
    from pptx.dml.color import RGBColor as PptRGB
    from pptx.enum.text import PP_ALIGN

    def _hex(h):
        h = h.strip("#")
        return PptRGB(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    bg_hex  = global_reqs.get("bg_color", "#1A1A1A")
    acc_hex = global_reqs.get("accent",   "#E91E8C")
    txt_hex = global_reqs.get("text_color", "#F2F2F2")
    BG1 = _hex(bg_hex)
    BG2 = _hex("2D2D2D")
    ACC = _hex(acc_hex)
    WHT = _hex(txt_hex)
    GRY = _hex("808080")
    GRN = _hex("2ECC71")
    RED = _hex("FF4444")
    YLW = _hex("FFCC00")

    brand   = product_info.get("brand", "")
    name    = product_info.get("name", "")
    article = product_info.get("article", "")
    FOOTER  = f"{brand}  |  {name}  |  Арт. {article}"

    prs = Presentation()
    prs.slide_width  = Cm(25.4)
    prs.slide_height = Cm(33.9)

    def blank(bg=BG1):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        f = s.background.fill
        f.solid()
        f.fore_color.rgb = bg
        return s

    def tb(slide, text, l, t, w, h, sz=12, bold=False, color=WHT,
           align=PP_ALIGN.LEFT, italic=False, wrap=True):
        box = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = PptPt(sz)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return box

    def rct(slide, l, t, w, h, fill=None, line=None, lw=PptPt(0.5)):
        sh = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
        else:
            sh.fill.background()
        if line:
            sh.line.color.rgb = line
            sh.line.width = lw
        else:
            sh.line.fill.background()
        return sh

    def hline(slide, top, left=0.5, right=24.9, color=None):
        c = color or ACC
        rct(slide, left, top, right - left, 0.07, fill=c)

    def foot(slide):
        tb(slide, FOOTER, 0.5, 33.0, 24.4, 0.8, sz=9, color=GRY)

    def sec_label(slide, text, top):
        tb(slide, text, 0.8, top, 23.9, 0.9, sz=15, bold=True, color=ACC)
        return top + 1.0

    def mbullets(slide, items, top, sz=12):
        """items: list of (sym_str, color, text)."""
        for sym, col, text in items:
            tb(slide, sym,  0.8, top, 0.9, 0.7, sz=sz, bold=True, color=col)
            tb(slide, text, 1.8, top, 22.5, 0.9, sz=sz, color=WHT)
            top += 0.95
        return top

    def two_col(slide, rows, top, sz=11):
        for label, value in rows:
            tb(slide, label, 0.8,  top, 8.0, 0.7, sz=sz, bold=True, color=ACC)
            tb(slide, value, 9.2,  top, 15.0, 0.7, sz=sz, color=WHT)
            top += 0.82
        return top

    # ── Слайд 1: Общие требования ─────────────────────────────────────────
    s1 = blank()
    tb(s1, "ОБЩИЕ ТРЕБОВАНИЯ", 0.5, 0.5, 24.4, 1.5, sz=28, bold=True, color=ACC)
    hline(s1, 2.2)

    top = 2.5
    top = sec_label(s1, "Формат и качество:", top)
    top = two_col(s1, [
        ("Формат карточек", global_reqs.get("formats", "1200×1600 px (инфографика) | 1200×1200 px (обложка)")),
        ("Качество",        "300 dpi, RGB"),
        ("Ориентация",      "Вертикальная (портрет)"),
        ("Формат сдачи",    "PNG/JPEG (финал) + PSD/AI (макет)"),
    ], top)
    top += 0.3

    hline(s1, top, color=GRY)
    top += 0.3
    top = sec_label(s1, "Цветовая палитра:", top)

    # Color swatches
    swatches = [
        (bg_hex,  "Основной фон"),
        (acc_hex, "Акцент бренда"),
        (txt_hex, "Текст основной"),
        ("#FFFFFF","Текст вторичный"),
    ]
    for i, (hx, lbl) in enumerate(swatches):
        lft = 0.8 + i * 6.0
        rct(s1, lft, top, 5.2, 1.8, fill=_hex(hx), line=GRY, lw=PptPt(0.3))
        tb(s1, hx,  lft, top + 1.9, 5.2, 0.6, sz=10, bold=True, color=ACC)
        tb(s1, lbl, lft, top + 2.6, 5.2, 0.6, sz=10, color=GRY)

    top += 3.4
    hline(s1, top, color=GRY)
    top += 0.3
    top = sec_label(s1, "Стиль и настроение:", top)
    palette_note = global_reqs.get("palette_note", "")
    if palette_note:
        tb(s1, palette_note, 0.8, top, 23.9, 1.2, sz=12, color=WHT)
        top += 1.5

    hline(s1, top, color=GRY)
    top += 0.3
    tb(s1, "ЗАПРЕЩЕНО:", 0.8, top, 23.9, 0.9, sz=15, bold=True, color=RED)
    top += 0.9
    forbidden = global_reqs.get("forbidden", [])
    top = mbullets(s1, [("✗", RED, f) for f in forbidden], top, sz=12)

    foot(s1)

    # ── Слайды карточек ───────────────────────────────────────────────────
    for card in cards:
        s = blank()
        num = card.get("number", "?")
        title = card.get("title", "")
        tb(s, str(num), 0.5, 0.4, 2.0, 2.0, sz=40, bold=True, color=ACC)
        tb(s, title, 2.8, 0.6, 21.6, 1.8, sz=20, bold=True, color=WHT)
        hline(s, 2.7)

        top = 3.0
        goal = card.get("goal", "")
        if goal:
            tb(s, f"Цель: {goal}", 0.8, top, 23.9, 0.8, sz=12, italic=True, color=_hex("FFA040"))
            top += 1.0

        bg_desc = card.get("bg", "")
        if bg_desc:
            top = sec_label(s, "Фон:", top)
            tb(s, bg_desc, 0.8, top, 23.9, 1.5, sz=12, color=WHT)
            top += 1.8

        comp_desc = card.get("composition", "")
        if comp_desc:
            top = sec_label(s, "Композиция:", top)
            tb(s, comp_desc, 0.8, top, 23.9, 2.0, sz=12, color=WHT)
            top += 2.3

        text_blocks = card.get("text_blocks", [])
        if text_blocks:
            top = sec_label(s, "Текстовые блоки:", top)
            top = two_col(s, text_blocks, top, sz=11)
            top += 0.3

        comp_ref = card.get("competitor_ref", "")
        if comp_ref:
            hline(s, top, color=GRY)
            top += 0.3
            top = sec_label(s, "Справка для дизайнера (конкуренты):", top)
            tb(s, comp_ref, 0.8, top, 23.9, 2.0, sz=11, color=GRY)
            top += 2.3

        do_dont = card.get("do_dont", [])
        if do_dont:
            rows = []
            for item in do_dont:
                if len(item) == 2:
                    sym, text = item
                    col = GRN if sym == "✓" else RED
                    rows.append((sym, col, text))
                else:
                    rows.append(item)
            mbullets(s, rows, top, sz=12)

        foot(s)

    # ── Технические требования ────────────────────────────────────────────
    s_tech = blank(BG2)
    tb(s_tech, "ТЕХНИЧЕСКИЕ ТРЕБОВАНИЯ", 0.5, 0.4, 24.4, 1.5, sz=28, bold=True, color=ACC)
    hline(s_tech, 2.2)

    top = 2.5
    top = sec_label(s_tech, "Шрифты:", top)
    for line in [
        "Заголовки: жирный, без засечек — Arial Bold / Montserrat Bold",
        "Основной текст: чёткий, контрастный — Arial / Roboto",
        "Минимальный размер: 24 px (читаемость на мобильном 375 px)",
        "Не использовать декоративные и курсивные шрифты для ключевых надписей",
    ]:
        tb(s_tech, f"• {line}", 0.8, top, 23.9, 0.7, sz=13, color=WHT)
        top += 0.75

    top += 0.3
    hline(s_tech, top, color=GRY)
    top += 0.3
    top = sec_label(s_tech, "Цветовая палитра:", top)

    palette_entries = [
        (bg_hex,  "Основной фон",    "Все слайды"),
        (acc_hex, "Акцент бренда",   "Заголовки, маркеры, рамки"),
        (txt_hex, "Текст основной",  "Все подписи"),
        ("#FFFFFF","Текст вторичный", "Описания, пояснения"),
        ("#FFCC00","Предупреждения",  "⚠ блоки"),
        ("#FF4444","Запрещено",       "✗ блоки"),
        ("#2ECC71","Обязательно",     "✓ блоки"),
    ]
    for hx, lbl, usage in palette_entries:
        rct(s_tech, 0.8, top, 1.2, 0.6, fill=_hex(hx), line=GRY, lw=PptPt(0.3))
        tb(s_tech, hx,    2.2,  top + 0.05, 3.5, 0.55, sz=10, bold=True, color=ACC)
        tb(s_tech, lbl,   6.0,  top + 0.05, 8.0, 0.55, sz=10, color=WHT)
        tb(s_tech, usage, 14.5, top + 0.05, 10.5, 0.55, sz=10, color=GRY)
        top += 0.68

    top += 0.3
    hline(s_tech, top, color=GRY)
    top += 0.3
    top = sec_label(s_tech, f"ВАЖНЫЕ АКЦЕНТЫ — УТП {brand}:", top)
    utp = global_reqs.get("utp_bullets", [])
    top = mbullets(s_tech, [("✓", GRN, u) for u in utp], top, sz=12)

    top += 0.3
    hline(s_tech, top, color=GRY)
    top += 0.3
    top = sec_label(s_tech, f"СОСТАВ ПАКЕТА ФОТОГРАФИЙ (итого {len(cards)} изображений):", top)
    for i, card in enumerate(cards, 1):
        tb(s_tech, f"{i}. {card.get('title', '')}  — 1200×1600 px",
           0.8, top, 23.9, 0.7, sz=12, color=WHT)
        top += 0.72

    foot(s_tech)

    # ── Анализ конкурентов (таблица) ──────────────────────────────────────
    s_comp = blank(BG2)
    tb(s_comp, f"АНАЛИЗ КОНКУРЕНТОВ → ФОКУС ДЛЯ КАЖДОГО ФОТО {brand}",
       0.5, 0.3, 24.4, 1.5, sz=17, bold=True, color=ACC)
    tb(s_comp, "Что делают лидеры и на что обратить внимание при создании нашей карточки",
       0.5, 1.9, 24.4, 0.8, sz=11, color=GRY)
    hline(s_comp, 2.8)

    top = 3.0
    rct(s_comp, 0.5, top, 24.4, 0.9, fill=ACC)
    tb(s_comp, "№",      0.6,  top + 0.1, 0.7,  0.7, sz=11, bold=True, color=WHT)
    tb(s_comp, "Наше фото", 1.5, top + 0.1, 3.0, 0.7, sz=11, bold=True, color=WHT)
    tb(s_comp, "Что делают конкуренты (лучший пример)", 4.8, top + 0.1, 9.5, 0.7, sz=11, bold=True, color=WHT)
    tb(s_comp, f"На что обратить внимание для {brand}", 14.5, top + 0.1, 10.4, 0.7, sz=11, bold=True, color=WHT)
    top += 0.95

    rows_ct = competitor_table or []
    ROW_H = max(2.8, round((33.0 - top) / max(len(rows_ct), 1), 2))
    for i, row in enumerate(rows_ct):
        bg = BG1 if i % 2 == 0 else BG2
        rct(s_comp, 0.5, top, 24.4, ROW_H, fill=bg)
        for lft, wid in [(1.4, 0.05), (4.7, 0.05), (14.4, 0.05)]:
            rct(s_comp, lft, top, wid, ROW_H, fill=GRY)
        tb(s_comp, str(row.get("num", i+1)), 0.6, top + 0.1, 0.7, ROW_H - 0.2, sz=12, bold=True, color=ACC)
        tb(s_comp, row.get("photo", ""),     1.5, top + 0.1, 3.0, ROW_H - 0.2, sz=10, bold=True, color=WHT)
        tb(s_comp, row.get("competitors", ""), 4.8, top + 0.1, 9.5, ROW_H - 0.2, sz=9,  color=GRY)
        tb(s_comp, row.get("focus", ""),    14.5, top + 0.1, 10.4, ROW_H - 0.2, sz=9,  color=WHT)
        top += ROW_H

    foot(s_comp)

    # ── Слайды Фото 1..M: примеры конкурентов ────────────────────────────
    for psd in (photo_slides or []):
        s = blank()
        num_p = psd.get("photo_num", "?")
        rct(s, 0.5, 0.3, 24.4, 1.8, fill=BG2)
        tb(s, psd.get("title", f"ФОТО {num_p}: примеры конкурентов"),
           0.8, 0.4, 23.9, 1.5, sz=16, bold=True, color=ACC)
        hline(s, 2.2, color=GRY)

        top_ph = 2.5
        block_w = 5.8
        competitors_list = psd.get("competitors", [])
        for i, comp in enumerate(competitors_list[:4]):
            brand_name = comp[0] if len(comp) > 0 else "—"
            brand_desc = comp[1] if len(comp) > 1 else ""
            lft = 0.6 + i * (block_w + 0.4)
            rct(s, lft, top_ph, block_w, 12.5, fill=BG2, line=GRY, lw=PptPt(0.5))
            rct(s, lft + 2.5, top_ph + 0.3, 3.0, 0.8, fill=ACC)
            tb(s, brand_name, lft + 2.5, top_ph + 0.3, 3.0, 0.8,
               sz=10, bold=True, color=WHT, align=PP_ALIGN.CENTER)
            rct(s, lft + 0.2, top_ph + 1.3, block_w - 0.4, 7.5, fill=BG1, line=GRY, lw=PptPt(0.3))
            tb(s, "[ФОТО\nконкурента]", lft + 0.2, top_ph + 4.2, block_w - 0.4, 2.0,
               sz=11, color=GRY, align=PP_ALIGN.CENTER)
            tb(s, brand_desc, lft + 0.2, top_ph + 9.0, block_w - 0.4, 3.2, sz=10, color=WHT)

        top_n = top_ph + 13.1
        hline(s, top_n, color=GRY)
        top_n += 0.3
        tb(s, f"▶  НА ЧТО ОБРАТИТЬ ВНИМАНИЕ — ФОТО {num_p} {brand}",
           0.5, top_n, 24.4, 0.9, sz=13, bold=True, color=ACC)
        top_n += 1.1
        notes = psd.get("notes", [])
        rows_n = []
        for note in notes:
            if len(note) == 2:
                sym, text = note
                col = GRN if sym == "✓" else RED
                rows_n.append((sym, col, text))
            else:
                rows_n.append(note)
        mbullets(s, rows_n, top_n, sz=12)
        foot(s)

    path = folder / filename
    prs.save(str(path))
    return str(path)


def make_docx_photographer_tz(folder: Path, filename: str,
                               product_info: dict, blocks: list) -> str:
    """
    Создаёт DOCX ТЗ для фотографа.

    product_info: {"brand", "name", "article", "date"}
    blocks: [{
        "number": int, "name": str, "goal": str,
        "shots": [{"frame": str, "description": str, "angle": str, "notes": str}],
        "background": str, "lighting": str, "props": str
    }]
    """
    doc = Document()
    for sec in doc.sections:
        sec.top_margin    = Cm(1.5)
        sec.bottom_margin = Cm(1.5)
        sec.left_margin   = Cm(2.0)
        sec.right_margin  = Cm(1.5)

    brand   = product_info.get("brand", "")
    name    = product_info.get("name", "")
    article = product_info.get("article", "")
    today   = datetime.date.today().strftime("%d.%m.%Y")

    # Заголовок
    h = doc.add_heading(f"ТЗ ДЛЯ ФОТОГРАФА — {brand} {name}", level=0)
    h.runs[0].font.color.rgb = DARK_RGB
    p = doc.add_paragraph(f"Дата: {today}  |  Арт. {article}")
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = GREY_RGB
    doc.add_paragraph("")

    # Шапка проекта (таблица)
    h2 = doc.add_heading("ИНФОРМАЦИЯ О ПРОЕКТЕ", level=1)
    for run in h2.runs:
        run.font.color.rgb = GOLD_RGB
    info_rows = [
        ("Бренд",    brand),
        ("Товар",    name),
        ("Артикул",  article),
        ("Дата ТЗ",  today),
        ("Блоков съёмки", str(len(blocks))),
    ]
    tbl = doc.add_table(rows=len(info_rows), cols=2)
    tbl.style = "Table Grid"
    tbl.columns[0].width = Cm(5)
    tbl.columns[1].width = Cm(12)
    for i, (f, v) in enumerate(info_rows):
        c0, c1 = tbl.cell(i, 0), tbl.cell(i, 1)
        c0.text = f
        c0.paragraphs[0].runs[0].font.bold = True
        c0.paragraphs[0].runs[0].font.size = Pt(10)
        _set_cell_bg(c0, "F5F5F5")
        c1.text = v
        c1.paragraphs[0].runs[0].font.size = Pt(10)
        _set_cell_bg(c1, WHITE_HEX if i % 2 == 0 else LIGHT_HEX)
    doc.add_paragraph("")

    # Сводная таблица блоков
    h2 = doc.add_heading("БЛОКИ СЪЁМКИ — ОБЗОР", level=1)
    for run in h2.runs:
        run.font.color.rgb = GOLD_RGB

    ov_headers = ["№", "Блок", "Цель", "Кадров"]
    tbl2 = doc.add_table(rows=1 + len(blocks), cols=4)
    tbl2.style = "Table Grid"
    for ci, h_text in enumerate(ov_headers):
        cell = tbl2.cell(0, ci)
        cell.text = h_text
        cell.paragraphs[0].runs[0].font.bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(10)
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _set_cell_bg(cell, DARK_HEX)
    for i, blk in enumerate(blocks, start=1):
        tbl2.cell(i, 0).text = str(blk.get("number", i))
        tbl2.cell(i, 1).text = blk.get("name", "")
        tbl2.cell(i, 2).text = blk.get("goal", "")
        tbl2.cell(i, 3).text = str(len(blk.get("shots", [])))
        for ci in range(4):
            tbl2.cell(i, ci).paragraphs[0].runs[0].font.size = Pt(10)
            _set_cell_bg(tbl2.cell(i, ci), WHITE_HEX if i % 2 == 1 else LIGHT_HEX)
    doc.add_paragraph("")

    # Детальное ТЗ по каждому блоку
    for blk in blocks:
        num  = blk.get("number", "")
        name = blk.get("name", "")
        h2 = doc.add_heading(f"БЛОК {num}: {name.upper()}", level=1)
        for run in h2.runs:
            run.font.color.rgb = GOLD_RGB

        # Общие параметры блока
        meta_rows = [
            ("Цель блока",   blk.get("goal", "")),
            ("Фон",          blk.get("background", "")),
            ("Освещение",    blk.get("lighting", "")),
            ("Реквизит",     blk.get("props", "")),
        ]
        tbl_m = doc.add_table(rows=len(meta_rows), cols=2)
        tbl_m.style = "Table Grid"
        tbl_m.columns[0].width = Cm(4)
        tbl_m.columns[1].width = Cm(13)
        for i, (f, v) in enumerate(meta_rows):
            c0, c1 = tbl_m.cell(i, 0), tbl_m.cell(i, 1)
            c0.text = f
            c0.paragraphs[0].runs[0].font.bold = True
            c0.paragraphs[0].runs[0].font.size = Pt(10)
            _set_cell_bg(c0, "F5F5F5")
            c1.text = v
            c1.paragraphs[0].runs[0].font.size = Pt(10)
        doc.add_paragraph("")

        # Таблица кадров
        shots = blk.get("shots", [])
        if shots:
            sh_headers = ["Кадр", "Описание", "Ракурс", "Примечания"]
            tbl_s = doc.add_table(rows=1 + len(shots), cols=4)
            tbl_s.style = "Table Grid"
            for ci, h_text in enumerate(sh_headers):
                cell = tbl_s.cell(0, ci)
                cell.text = h_text
                cell.paragraphs[0].runs[0].font.bold = True
                cell.paragraphs[0].runs[0].font.size = Pt(9)
                cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                _set_cell_bg(cell, GOLD_HEX)
            for i, shot in enumerate(shots, start=1):
                tbl_s.cell(i, 0).text = shot.get("frame", str(i))
                tbl_s.cell(i, 1).text = shot.get("description", "")
                tbl_s.cell(i, 2).text = shot.get("angle", "")
                tbl_s.cell(i, 3).text = shot.get("notes", "")
                for ci in range(4):
                    tbl_s.cell(i, ci).paragraphs[0].runs[0].font.size = Pt(9)
                    _set_cell_bg(tbl_s.cell(i, ci), WHITE_HEX if i % 2 == 1 else LIGHT_HEX)
        doc.add_paragraph("")

    # Требования к сдаче
    h2 = doc.add_heading("ТРЕБОВАНИЯ К СДАЧЕ", level=1)
    for run in h2.runs:
        run.font.color.rgb = GOLD_RGB
    deliverables = [
        "Форматы файлов: TIFF (мастер) + JPEG (веб) + RAW (архив)",
        "Разрешение: минимум 3000×4000 px (1200×1600 @ 300dpi)",
        "Цветовой профиль: RGB sRGB",
        "Именование: [Артикул]_Блок[N]_Кадр[N]_v1.tiff",
        "Ретушь: удаление пыли, коррекция баланса белого, без избыточной обработки",
        "Сдача: папка с подпапками по блокам + сводный JPEG-превью",
    ]
    for line in deliverables:
        p = doc.add_paragraph(line, style="List Bullet")
        if p.runs:
            p.runs[0].font.size = Pt(10)

    path = folder / filename
    doc.save(str(path))
    return str(path)


if __name__ == "__main__":
    folder = get_product_folder("ifoam", "OrangeStandard", "119105")
    print("Папка:", folder)

    path = make_docx_card(folder, "01_card_119105.docx", "Карточка товара — ifoam Orange Standard", [
        {
            "heading": "КАРТОЧКА ТОВАРА",
            "type": "table",
            "rows": [
                ("Бренд", "ifoam AUTO"),
                ("Название продукта", "Orange Standard"),
                ("Артикул / SKU", "Art. 119105"),
                ("Объём / масса", "5 кг"),
                ("pH", "12"),
                ("Химическая основа", "Щелочной"),
                ("Себестоимость", "496 ₽/шт"),
                ("Цена закупки", "530 ₽/шт"),
                ("В коробе", "4 шт"),
            ]
        },
        {
            "heading": "Аудитория и УТП",
            "type": "text",
            "text": "Товар ориентирован на профессиональные автомойки и продвинутых автолюбителей. "
                    "Главное УТП — концентрат pH 12 в объёме 5 кг: 1 л = 20 автомобилей."
        }
    ])
    print("DOCX карточка:", path)

    path = make_xlsx_table(
        folder, "test_table.xlsx", "Тест таблицы",
        ["Колонка 1", "Колонка 2"],
        [["Значение 1", "Значение 2"], ["Значение 3", "Значение 4"]]
    )
    print("XLSX:", path)
