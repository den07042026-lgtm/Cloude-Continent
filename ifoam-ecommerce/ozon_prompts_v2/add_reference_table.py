# -*- coding: utf-8 -*-
"""Добавляет слайд-таблицу «Анализ конкурентов → фокус для ifoam» в существующий PPTX"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
ORANGE     = RGBColor(0xFF, 0x6B, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY  = RGBColor(0x2D, 0x2D, 0x2D)
ORANGE_LITE= RGBColor(0xFF, 0xA0, 0x40)
GREEN_OK   = RGBColor(0x2E, 0xCC, 0x71)
RED_WARN   = RGBColor(0xFF, 0x44, 0x44)
YELLOW     = RGBColor(0xFF, 0xCC, 0x00)
ROW_ODD    = RGBColor(0x22, 0x22, 0x22)
ROW_EVEN   = RGBColor(0x1A, 0x1A, 0x1A)
HDR_BG     = RGBColor(0x33, 0x18, 0x00)

import sys
sys.path.insert(0, "C:/Users/1/Downloads/ozon_prompts_v2")
from save_helpers import get_product_folder
folder = get_product_folder("ifoam", "OrangeStandard5kg", "119105")
pptx_path = folder / "07_designer_tz_119105_v2.pptx"

prs = Presentation(str(pptx_path))
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          font_size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txb

def rect(slide, left, top, width, height, fill_color, line_color=None, lw=0):
    from pptx.util import Pt as PPt
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = PPt(lw)
    else:
        shape.line.fill.background()
    return shape

def divider(slide, top, color=ORANGE):
    rect(slide, 0.2, top, 9.6, 0.04, fill_color=color)

# ─── Данные таблицы ───────────────────────────────────────────
# Столбцы: №, Наше фото, Что делают конкуренты (лучший пример), На что обратить внимание для ifoam
rows = [
    (
        "1",
        "Главное фото\n(продакт-шот)",
        "GLITTER G10: чистый продакт-шот, светлый фон, этикетка чётко читается, минимум текста на фото\n"
        "GRASS Balance: яркий жёлтый канистр + «Бренд №1» бейдж, пена снизу\n"
        "DR.BERG ACE: тёмный фон, канистр крупно, нет лишнего",
        "✓ Тёмный фон (#1A1A1A) — отстройка от конкурентов\n"
        "✓ Оранжевый акцент — цвет жидкости виден через канистр\n"
        "✓ Этикетка «Orange Standard» читается полностью\n"
        "✓ «5 кг» бейдж + «ОТ ПРОИЗВОДИТЕЛЯ» метка\n"
        "✗ Не светлый фон — теряемся рядом с GLITTER G10",
    ),
    (
        "2",
        "Характеристики\n/ УТП",
        "GLITTER G10: 4 буллета с иконками — работает с водой / густая пена / смывается / не оставляет разводов\n"
        "DR.BERG ACE: 3 плашки на тёмном — «активная пена / плотная пена / безопасно для покрытия»\n"
        "GLITTER G5: «густая активная пена — удаляет грязь и реагенты / работает без губки»",
        "✓ pH 12 — крупно в отдельном блоке (конкуренты это НЕ делают)\n"
        "✓ Пропорции разбавления 1:30–1:40 / 1:3–1:4 — таблица или иконки\n"
        "✓ «Не повреждает ЛКП» — закрываем главное возражение\n"
        "✗ Не копировать зелёные маркеры — цвет GRASS\n"
        "✗ Не перегружать — не более 5–6 буллетов",
    ),
    (
        "3",
        "Инструкция по\nприменению",
        "GRASS Balance: пошаговые фото «1-откройте / 2-разведите / 3-остудите кузов» — тёмный стиль, реальные руки\n"
        "GLITTER G5: «4 шага до чистого кузова» — сетка 2×2 с фото процесса\n"
        "GLITTER G10: инструкция с иконками пеногенератор/пенокомплект/дозатрон + пропорции",
        "✓ 4 шага — остудите / разведите / нанесите / смойте\n"
        "✓ Шаг 2: КРУПНО пропорции — пеногенератор 1:30–1:40, пенокомплект 1:3–1:4\n"
        "✓ Реальные фото процесса (не иконки) — как у GRASS\n"
        "✓ Предупреждение «НЕ ДОПУСКАЙТЕ ВЫСЫХАНИЯ» — жёлтым\n"
        "✗ Не мелкий текст — пропорции должны читаться с телефона",
    ),
    (
        "4",
        "Экономика\n(уникальный слайд)",
        "DR.BERG ACE: «до 200 моек» на главном фото мелко\n"
        "BIGHIM: «5л = до 150 моек» на карточке характеристик\n"
        "Остальные конкуренты — ОТДЕЛЬНОГО слайда по экономике НЕТ ни у кого",
        "✓ Это наше УТП — делаем ПЕРВЫЙ в нише «экономический» слайд\n"
        "✓ Заголовок: «5 КГ = 20 АВТОМОБИЛЕЙ» или «МОЙКА ОТ 12 РУБЛЕЙ»\n"
        "✓ Три цифровых блока: 5 кг → 150–200 л → ~12 руб/авто\n"
        "✓ Таблица сравнения: автомойка 300–600р vs ifoam ~12р\n"
        "✗ Не перегружать расчётами — только 3 ключевые цифры крупно",
    ),
    (
        "5",
        "Рекомендации\n/ Безопасность",
        "GLITTER G10: «Наденьте перчатки» — крупно, фото синей перчатки на руке, серый фон\n"
        "GLITTER G5: «Защитные перчатки» — фото мужчины с пистолетом на мойке\n"
        "DR.BERG: предупреждения текстом — «не наносить на горячее / не допускать высыхания»",
        "✓ Реальное фото человека в перчатках (не иконка) — как у GLITTER\n"
        "✓ Главное предупреждение крупно: «НЕ НАНОСИТЬ НА ГОРЯЧИЙ КУЗОВ»\n"
        "✓ Второе: «ВЫДЕРЖАТЬ НЕ БОЛЕЕ 2 МИНУТ» — жёлтый цвет\n"
        "✓ Блок антивозражений: «что будет если...» — закрываем страхи\n"
        "✗ Не просто список — нужна визуальная иерархия приоритетов",
    ),
    (
        "6",
        "Производство\n/ Доверие",
        "GLITTER G5: «Сделано в России для наших дорог» — коллаж лаборатории 2×2, флаг РФ\n"
        "GLITTER G5: «51 000+ положительных отзывов» — карточки с цитатами и звёздами\n"
        "DR.BERG: «Немецкое качество» на главном фото, логотип DE\n"
        "GLITTER G5: «Мойка своими руками — как у профи» — фото мужчины с канистром",
        "✓ «ПРОИЗВОДИМ САМИ — ОТВЕЧАЕМ ЗА КАЧЕСТВО» — фото цеха ifoam\n"
        "✓ Флаг РФ + ТУ 20.41.32-003-36603872-2024 — закрываем «подделка?»\n"
        "✓ Коллаж 2×2: оборудование / лаборатория / розлив / готовая продукция\n"
        "✓ Если есть отзывы — 2 карточки с цитатами и ★★★★★\n"
        "✗ Не AI-фото лаборатории — только реальные кадры производства",
    ),
]

# ─── СЛАЙД 9 — ТАБЛИЦА ────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
bg(s9, DARK_GRAY)

txbox(s9, "АНАЛИЗ КОНКУРЕНТОВ → ФОКУС ДЛЯ КАЖДОГО ФОТО ifoam", 0.2, 0.12, 9.6, 0.7,
      font_size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txbox(s9, "Что делают лидеры и на что обратить внимание при создании нашей карточки",
      0.2, 0.78, 9.6, 0.38, font_size=13, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
divider(s9, 1.2)

# Заголовок таблицы
HDR_H = 0.45
col_w = [0.35, 1.25, 3.6, 4.0]
col_x = [0.2, 0.55, 1.8, 5.4]
headers = ["№", "Наше фото", "Что делают конкуренты (лучший пример)", "На что обратить внимание для ifoam"]

rect(s9, 0.2, 1.25, 9.6, HDR_H, fill_color=HDR_BG)
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    txbox(s9, hdr, cx + 0.05, 1.28, cw - 0.1, HDR_H - 0.05,
          font_size=13, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

# Строки
ROW_H = 1.88
start_y = 1.7

for i, (num, photo, competitor, focus) in enumerate(rows):
    top = start_y + i * ROW_H
    fill = ROW_ODD if i % 2 == 0 else ROW_EVEN
    rect(s9, 0.2, top, 9.6, ROW_H - 0.05, fill_color=fill,
         line_color=RGBColor(0x40,0x40,0x40), lw=0.5)

    # Вертикальные разделители
    for cx in [0.55, 1.8, 5.4]:
        rect(s9, cx, top, 0.02, ROW_H - 0.05, fill_color=RGBColor(0x40,0x40,0x40))

    # №
    txbox(s9, num, col_x[0] + 0.05, top + 0.05, col_w[0] - 0.1, ROW_H - 0.15,
          font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Наше фото
    txbox(s9, photo, col_x[1] + 0.06, top + 0.08, col_w[1] - 0.12, ROW_H - 0.15,
          font_size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Что делают конкуренты
    txbox(s9, competitor, col_x[2] + 0.06, top + 0.06, col_w[2] - 0.12, ROW_H - 0.12,
          font_size=11, color=LIGHT_GRAY, italic=False)

    # Фокус для ifoam — зелёные ✓ красные ✗
    txbox(s9, focus, col_x[3] + 0.06, top + 0.06, col_w[3] - 0.12, ROW_H - 0.12,
          font_size=11, color=WHITE)

# Нижняя легенда
divider(s9, 12.7)
txbox(s9, "✓ — делать обязательно     ✗ — избегать     → — уникальное преимущество ifoam",
      0.2, 12.78, 9.6, 0.38, font_size=12,
      color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.CENTER)

# ─── Сохранить ────────────────────────────────────────────────
prs.save(str(pptx_path))
print(f"Обновлено: {pptx_path}")
print(f"Слайдов: {len(prs.slides)}")
