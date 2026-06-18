# -*- coding: utf-8 -*-
"""Перерисовывает слайд 5 (Экономика ТЗ) — без рублей, только расход и объём."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

sys.path.insert(0, "C:/Users/1/Downloads/ozon_prompts_v2")
from save_helpers import get_product_folder
folder    = get_product_folder("ifoam", "OrangeStandard5kg", "119105")
pptx_path = folder / "07_designer_tz_119105_v2.pptx"

prs = Presentation(str(pptx_path))

# ── Очистить все shapes на слайде 5 (индекс 4) ───────────────
s = prs.slides[4]
spTree = s.shapes._spTree
# удаляем все дочерние элементы кроме spGrpSp (placeholder группы)
for el in list(spTree):
    tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
    if tag not in ('grpSpPr', 'grpSp'):
        spTree.remove(el)

# ── цвета ─────────────────────────────────────────────────────
BLACK      = RGBColor(0x1A,0x1A,0x1A)
ORANGE     = RGBColor(0xFF,0x6B,0x00)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
LIGHT_GRAY = RGBColor(0xF2,0xF2,0xF2)
DARK_GRAY  = RGBColor(0x2D,0x2D,0x2D)
ORANGE_LITE= RGBColor(0xFF,0xA0,0x40)
GREEN      = RGBColor(0x2E,0xCC,0x71)
RED        = RGBColor(0xFF,0x44,0x44)
YELLOW     = RGBColor(0xFF,0xCC,0x00)
HDR_BG     = RGBColor(0x33,0x18,0x00)

def txbox(slide, text, left, top, width, height,
          font_size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txb

def rect(slide, left, top, width, height, fill_color, line_color=None, lw=0):
    from pptx.util import Pt as PPt
    sh = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = PPt(lw)
    else:
        sh.line.fill.background()
    return sh

def divider(top):
    rect(s, 0.2, top, 9.6, 0.04, fill_color=ORANGE)

# ── фон ──────────────────────────────────────────────────────
fill = s.background.fill
fill.solid()
fill.fore_color.rgb = BLACK

# ════════════════════════════════════════════════════════════════
# ЗАГОЛОВОК РАЗДЕЛА
# ════════════════════════════════════════════════════════════════
txbox(s, "4", 0.3, 0.25, 0.7, 0.7,
      font_size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txbox(s, 'ФОТО "ЭКОНОМИКА РАСХОДА"', 1.1, 0.3, 8.5, 0.65,
      font_size=28, bold=True, color=WHITE)
divider(1.05)

txbox(s, "Концепция: показываем сколько выходит из одной канистры — без цен",
      0.3, 1.18, 9.4, 0.45, font_size=16, italic=True, color=ORANGE_LITE)

# ════════════════════════════════════════════════════════════════
# ВИЗУАЛЬНАЯ ЦЕПОЧКА (главный элемент слайда)
# ════════════════════════════════════════════════════════════════
txbox(s, "Заголовок слайда (крупно, центр):", 0.3, 1.75, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
txbox(s, '"5 КГ = 20 АВТОМОБИЛЕЙ"  или  "1 КГ КОНЦЕНТРАТА = 30–40 ЛИТРОВ РАСТВОРА"',
      0.3, 2.22, 9.4, 0.55, font_size=17, bold=True, color=WHITE)

divider(2.9)

# Три блока-стрелки
txbox(s, "Визуальная цепочка (3 блока + стрелки):", 0.3, 3.05, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)

chain = [
    ("🪣", "5 кг\nконцентрата",         "Исходный объём\nупаковки"),
    ("💧", "150–200 л\nрабочего раствора", "При разведении\n1:30–1:40"),
    ("🚗", "≥ 20\nавтомобилей",           "Количество\nполных моек"),
]
BW = 2.65
for i, (icon, val, sub) in enumerate(chain):
    bx = 0.3 + i*(BW + 0.6)
    rect(s, bx, 3.6, BW, 2.1,
         fill_color=RGBColor(0x2A,0x14,0x00),
         line_color=ORANGE, lw=2)
    txbox(s, icon, bx, 3.65, BW, 0.55,
          font_size=28, align=PP_ALIGN.CENTER)
    txbox(s, val, bx, 4.15, BW, 0.75,
          font_size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txbox(s, sub, bx, 4.9, BW, 0.65,
          font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)
    if i < 2:
        txbox(s, "→", bx+BW+0.1, 4.2, 0.42, 0.65,
              font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider(5.88)

# ════════════════════════════════════════════════════════════════
# ТАБЛИЦА ПРОПОРЦИЙ РАЗБАВЛЕНИЯ
# ════════════════════════════════════════════════════════════════
txbox(s, "Таблица пропорций разбавления:", 0.3, 6.05, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)

# Заголовок таблицы
rect(s, 0.3, 6.55, 9.4, 0.5, fill_color=HDR_BG)
txbox(s, "Способ нанесения", 0.4, 6.6, 3.5, 0.4,
      font_size=15, bold=True, color=ORANGE)
txbox(s, "Пропорция", 3.95, 6.6, 2.8, 0.4,
      font_size=15, bold=True, color=ORANGE)
txbox(s, "Расход на 1 л воды", 6.85, 6.6, 2.8, 0.4,
      font_size=15, bold=True, color=ORANGE)

# Строки
rows = [
    ("Пеногенератор (профи / автомойка)", "1:30 – 1:40", "20–25 г"),
    ("Пенокомплект (ручная мойка дома)",   "1:3 – 1:4",  "200–250 г"),
]
for j, (c1, c2, c3) in enumerate(rows):
    top = 7.1 + j * 0.6
    bg_c = RGBColor(0x28,0x28,0x28) if j%2==0 else RGBColor(0x22,0x22,0x22)
    rect(s, 0.3, top, 9.4, 0.55, fill_color=bg_c)
    txbox(s, c1, 0.4, top+0.07, 3.5, 0.42, font_size=15, color=LIGHT_GRAY)
    txbox(s, c2, 3.95, top+0.07, 2.8, 0.42, font_size=15, bold=True, color=ORANGE)
    txbox(s, c3, 6.85, top+0.07, 2.8, 0.42, font_size=15, color=WHITE)

divider(8.42)

# ════════════════════════════════════════════════════════════════
# ДОПОЛНИТЕЛЬНЫЕ ЭЛЕМЕНТЫ
# ════════════════════════════════════════════════════════════════
txbox(s, "Дополнительные элементы слайда:", 0.3, 8.58, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)

bullets = [
    ("✓", "Иконка канистры → стрелка → ведро → стрелка → автомобиль (инфографика)"),
    ("✓", "Подпись: «30–50 г концентрата на одно авто» — конкретная цифра расхода"),
    ("✓", "Акцент: «Концентрат — не значит сложно: просто добавь воды»"),
    ("✓", "Фоновое фото: автомобиль в густой пене или пеногенератор в действии"),
    ("✗", "Без цен, рублей, сравнений стоимости с конкурентами"),
    ("✗", "Без фраз «выгодно», «дёшево» — только факты расхода и объёма"),
]
for j, (mark, text) in enumerate(bullets):
    mc = GREEN if mark == "✓" else RED
    y  = 9.08 + j * 0.42
    txbox(s, mark, 0.3, y, 0.4, 0.4,
          font_size=14, bold=True, color=mc, align=PP_ALIGN.CENTER)
    txbox(s, text, 0.75, y, 9.0, 0.4, font_size=14, color=WHITE)

txbox(s, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)

# ── сохранить ─────────────────────────────────────────────────
prs.save(str(pptx_path))
print(f"Слайд 5 обновлён: {pptx_path}  |  Слайдов: {len(prs.slides)}")
