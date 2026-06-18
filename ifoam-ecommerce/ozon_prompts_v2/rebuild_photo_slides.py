# -*- coding: utf-8 -*-
"""
Удаляет старые слайды 10-15 из PPTX и добавляет новые с исправленным слайдом Экономика
(без рублей — только расход и концентрация).
"""
import sys, io, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, "C:/Users/1/Downloads/ozon_prompts_v2")
from save_helpers import get_product_folder
folder    = get_product_folder("ifoam", "OrangeStandard5kg", "119105")
pptx_path = folder / "07_designer_tz_119105_v2.pptx"
photos_dir = Path("C:/Users/1/Downloads/ifoam_OrangeStandard5kg_119105/Фото карточек конкурентов")
tmp_dir   = Path("C:/Users/1/Downloads/_tmp_png3")
tmp_dir.mkdir(exist_ok=True)

def to_png(stem):
    src = photos_dir / f"{stem}.webp"
    dst = tmp_dir / f"{stem}.png"
    if not dst.exists():
        Image.open(src).convert("RGB").save(dst, "PNG")
    return str(dst)

# ── цвета ─────────────────────────────────────────────────────
BLACK      = RGBColor(0x1A,0x1A,0x1A)
ORANGE     = RGBColor(0xFF,0x6B,0x00)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
LIGHT_GRAY = RGBColor(0xF2,0xF2,0xF2)
HDR_BG     = RGBColor(0x33,0x18,0x00)
GREEN      = RGBColor(0x2E,0xCC,0x71)
RED        = RGBColor(0xFF,0x44,0x44)
ORANGE_LITE= RGBColor(0xFF,0xA0,0x40)
BRAND_CLR  = {
    "GLITTER": RGBColor(0x00,0xAA,0xDD),
    "DR.BERG": RGBColor(0x66,0xCC,0x22),
    "GRASS":   RGBColor(0x22,0xAA,0x33),
    "BIGHIM":  RGBColor(0xFF,0x88,0x00),
}

prs   = Presentation(str(pptx_path))
BLANK = prs.slide_layouts[6]

# ── удалить все слайды начиная с индекса 9 (слайд 10 и далее) ─
sldIdLst = prs.slides._sldIdLst
while len(sldIdLst) > 9:
    sldIdLst.remove(sldIdLst[9])
print(f"После удаления: {len(prs.slides)} слайдов")

# ════════════════════════════════════════════════════════════════
# Вспомогательные функции
# ════════════════════════════════════════════════════════════════
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          font_size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = "Arial"
    return txb

def rect(slide, left, top, width, height, fill_color, line_color=None, lw=0):
    from pptx.util import Pt as PPt
    sh = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = PPt(lw)
    else:
        sh.line.fill.background()
    return sh

def divider(slide, top):
    rect(slide, 0.2, top, 9.6, 0.04, fill_color=ORANGE)

def add_photo(slide, png_path, left, top, width, height):
    img = Image.open(png_path)
    iw, ih = img.size
    ratio     = iw / ih
    box_ratio = width / height
    if ratio > box_ratio:
        w, h = width, width / ratio
        l, t = left, top + (height - h) / 2
    else:
        h, w = height, height * ratio
        l, t = left + (width - w) / 2, top
    slide.shapes.add_picture(png_path, Inches(l), Inches(t), Inches(w), Inches(h))

def footer(slide):
    txbox(slide, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
          0.2, 12.9, 9.6, 0.35, font_size=11,
          color=RGBColor(0x66,0x66,0x66), align=PP_ALIGN.CENTER)

def make_photo_slide(title, photos, focus_items, cols=4):
    """Создаёт слайд с фото конкурентов + блоком фокуса."""
    s = prs.slides.add_slide(BLANK)
    bg(s, BLACK)

    # Заголовок
    rect(s, 0.2, 0.12, 9.6, 0.72, fill_color=HDR_BG)
    txbox(s, title, 0.3, 0.17, 9.4, 0.65,
          font_size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    divider(s, 0.88)

    n      = len(photos)
    cols   = min(n, cols)
    MARGIN = 0.22
    GAP    = 0.15
    PH     = 3.8     # высота фото
    CAP_H  = 0.85    # высота подписи
    PT     = 1.0     # top фото

    photo_w = (10.0 - 2*MARGIN - GAP*(cols-1)) / cols

    for idx, (stem, caption, brand) in enumerate(photos):
        col  = idx % cols
        row  = idx // cols
        left = MARGIN + col*(photo_w + GAP)
        top  = PT + row*(PH + CAP_H + 0.1)

        png = to_png(stem)
        rect(s, left, top, photo_w, PH+CAP_H+0.08,
             fill_color=RGBColor(0x22,0x22,0x22),
             line_color=ORANGE, lw=1.0)
        add_photo(s, png, left+0.05, top+0.05, photo_w-0.1, PH-0.1)

        bc = BRAND_CLR.get(brand, ORANGE)
        rect(s, left+photo_w-1.3, top+0.05, 1.22, 0.32, fill_color=bc)
        txbox(s, brand, left+photo_w-1.3, top+0.05, 1.22, 0.32,
              font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(s, caption, left+0.06, top+PH+0.02, photo_w-0.12, CAP_H,
              font_size=11, color=LIGHT_GRAY)

    used_rows = (n + cols - 1) // cols
    focus_top = PT + used_rows*(PH + CAP_H + 0.1) + 0.1

    divider(s, focus_top)
    num = title.split(":")[0].replace("ФОТО","").strip()
    txbox(s, f"▶  НА ЧТО ОБРАТИТЬ ВНИМАНИЕ — ФОТО {num} ifoam Orange Standard",
          0.2, focus_top+0.1, 9.6, 0.45,
          font_size=15, bold=True, color=ORANGE)

    available_h = 13.33 - focus_top - 0.65
    line_h = min(0.42, available_h / max(len(focus_items), 1))

    for j, (mark, text) in enumerate(focus_items):
        mc = GREEN if mark == "✓" else (RED if mark == "✗" else ORANGE_LITE)
        y  = focus_top + 0.6 + j*line_h
        txbox(s, mark, 0.2, y, 0.4, line_h,
              font_size=14, bold=True, color=mc, align=PP_ALIGN.CENTER)
        txbox(s, text, 0.65, y, 9.1, line_h, font_size=13, color=WHITE)

    footer(s)
    return s

# ════════════════════════════════════════════════════════════════
# ДАННЫЕ СЛАЙДОВ
# ════════════════════════════════════════════════════════════════

# ── Слайд 10: ГЛАВНОЕ ФОТО ────────────────────────────────────
make_photo_slide(
    title="ФОТО 1: ГЛАВНОЕ ФОТО — примеры конкурентов",
    photos=[
        ("6592383288", "GLITTER G10\nЧистый продакт-шот,\nсветлый фон, этикетка читается",       "GLITTER"),
        ("8278266808", "DR.BERG ACE\nТёмный фон, канистр крупно,\n«до 200 моек» в кадре",         "DR.BERG"),
        ("9463559615", "GRASS Balance\n«Бренд №1» бейдж, яркий\nжёлтый — читается сразу",         "GRASS"),
        ("8901708053", "BIGHIM\n«5л = до 150 моек» крупно,\nяркий оранжевый фон",                 "BIGHIM"),
    ],
    focus_items=[
        ("✓", "Тёмный фон (#1A1A1A) — отстройка от светлого GLITTER G10"),
        ("✓", "Оранжевый акцент — цвет жидкости виден через канистру"),
        ("✓", "Этикетка Orange Standard читается полностью"),
        ("✓", "Бейдж «5 кг» + метка «ОТ ПРОИЗВОДИТЕЛЯ»"),
        ("✗", "Не белый фон — на полке сливаемся с GLITTER G10"),
    ],
)

# ── Слайд 11: ХАРАКТЕРИСТИКИ ──────────────────────────────────
make_photo_slide(
    title="ФОТО 2: ХАРАКТЕРИСТИКИ / УТП — примеры конкурентов",
    photos=[
        ("7388654580", "GLITTER G10\n4 буллета с иконками:\nпена / вода / смывается / разводы",    "GLITTER"),
        ("8278219581", "DR.BERG ACE\n3 плашки на тёмном:\nактивная пена / плотная / безопасно",    "DR.BERG"),
        ("9405795810", "GLITTER G5\n«Готов к применению»\nбуллеты + логотип",                     "GLITTER"),
        ("8901704980", "BIGHIM\nСписок на ярком фоне:\nкислотно-щелочной / блеск / водоотталк.",   "BIGHIM"),
    ],
    focus_items=[
        ("✓", "pH 12 — крупный блок отдельно (никто из конкурентов так НЕ делает)"),
        ("✓", "Пропорции 1:30–1:40 / 1:3–1:4 — иконки или таблица"),
        ("✓", "«Не повреждает ЛКП» — закрываем главное возражение"),
        ("✗", "Не зелёные маркеры — цвет занят брендом GRASS"),
        ("✗", "Не более 5–6 буллетов — не перегружать"),
    ],
)

# ── Слайд 12: ИНСТРУКЦИЯ ──────────────────────────────────────
make_photo_slide(
    title="ФОТО 3: ИНСТРУКЦИЯ ПО ПРИМЕНЕНИЮ — примеры конкурентов",
    photos=[
        ("7448143493", "GLITTER G10\nИнструкция с иконками:\nпеногенератор / пенокомплект / дозатрон", "GLITTER"),
        ("9405795764", "GLITTER G5\n«4 шага до чистого кузова»\nСетка 2×2 с реальными фото",           "GLITTER"),
        ("9463559637", "GRASS Balance\nШаг 1: «Откройте упаковку»\nРеальные руки, тёмный стиль",       "GRASS"),
        ("9463559647", "GRASS Balance\nШаг 2: «Разведите состав»\nПоказан концентрат наглядно",        "GRASS"),
    ],
    focus_items=[
        ("✓", "4 шага: остудите / разведите / нанесите / смойте"),
        ("✓", "Шаг 2: КРУПНО пропорции — 1:30–1:40 и 1:3–1:4"),
        ("✓", "Реальные фото процесса (не иконки) — как у GRASS"),
        ("✓", "Предупреждение «НЕ ДОПУСКАЙТЕ ВЫСЫХАНИЯ» жёлтым"),
        ("✗", "Не мелкий текст — пропорции должны читаться с телефона"),
    ],
)

# ── Слайд 13: ЭКОНОМИКА РАСХОДА (без рублей) ─────────────────
s13 = prs.slides.add_slide(BLANK)
bg(s13, BLACK)

rect(s13, 0.2, 0.12, 9.6, 0.72, fill_color=HDR_BG)
txbox(s13, "ФОТО 4: ЭКОНОМИКА РАСХОДА — как конкуренты показывают концентрацию",
      0.3, 0.17, 9.4, 0.65, font_size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
divider(s13, 0.88)

# Фото 2×2
eco_photos = [
    ("8278219912", "DR.BERG ACE\nТаблица разбавления — пеногенератор\n1:100–200, пенокомплект 1:10–20", "DR.BERG"),
    ("8278220387", "DR.BERG ACE\n«Высокая концентрация средства»\n+ «эффективен при жёсткой воде»",     "DR.BERG"),
    ("8901708053", "BIGHIM\n«5л = до 150 моек» крупно на фото\n— цифра выходит вперёд",                 "BIGHIM"),
    ("8278266808", "DR.BERG ACE\n«Автошампунь до 200 моек»\nна главном продакт-шоте",                   "DR.BERG"),
]
cols, MARGIN, GAP = 2, 0.22, 0.18
PH, CAP_H, PT = 3.5, 0.82, 1.0
photo_w = (10.0 - 2*MARGIN - GAP*(cols-1)) / cols

for idx, (stem, caption, brand) in enumerate(eco_photos):
    col  = idx % cols
    row  = idx // cols
    left = MARGIN + col*(photo_w+GAP)
    top  = PT + row*(PH+CAP_H+0.1)
    png  = to_png(stem)
    rect(s13, left, top, photo_w, PH+CAP_H+0.08,
         fill_color=RGBColor(0x22,0x22,0x22), line_color=ORANGE, lw=1.0)
    add_photo(s13, png, left+0.05, top+0.05, photo_w-0.1, PH-0.1)
    bc = BRAND_CLR.get(brand, ORANGE)
    rect(s13, left+photo_w-1.3, top+0.05, 1.22, 0.32, fill_color=bc)
    txbox(s13, brand, left+photo_w-1.3, top+0.05, 1.22, 0.32,
          font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s13, caption, left+0.06, top+PH+0.02, photo_w-0.12, CAP_H,
          font_size=11, color=LIGHT_GRAY)

# Блок фокуса
focus_top = PT + 2*(PH+CAP_H+0.1) + 0.05
divider(s13, focus_top)
txbox(s13, "▶  НА ЧТО ОБРАТИТЬ ВНИМАНИЕ — ФОТО 4 ifoam Orange Standard",
      0.2, focus_top+0.1, 9.6, 0.45, font_size=15, bold=True, color=ORANGE)

# Цепочка: канистра → литры → авто
BW = 2.75
for i, (lbl, arrow) in enumerate([
    ("5 кг\nконцентрата", "→"),
    ("150–200 л\nрабочего раствора", "→"),
    ("≥ 20\nавтомобилей", ""),
]):
    bx = 0.25 + i*(BW+0.58)
    rect(s13, bx, focus_top+0.62, BW, 1.05,
         fill_color=RGBColor(0x2A,0x14,0x00), line_color=ORANGE, lw=2)
    txbox(s13, lbl, bx, focus_top+0.67, BW, 0.95,
          font_size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    if arrow:
        txbox(s13, arrow, bx+BW+0.07, focus_top+0.85, 0.42, 0.6,
              font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

eco_focus = [
    ("✓", "Заголовок: «5 КГ = 20 АВТОМОБИЛЕЙ» — главная цифра вперёд"),
    ("✓", "Визуальная цепочка: канистра → ведро раствора → автомобили"),
    ("✓", "Пропорции разбавления крупно: пеногенератор 1:30–1:40 / пенокомплект 1:3–1:4"),
    ("✓", "Подпись: «1 кг концентрата = 30–40 литров рабочего состава»"),
    ("✗", "Без цен и рублей — только объём, пропорции, количество моек"),
    ("✗", "Не сравнивать с конкурентами — только наши показатели расхода"),
]
for j, (mark, text) in enumerate(eco_focus):
    mc = GREEN if mark == "✓" else RED
    y  = focus_top + 1.82 + j*0.38
    txbox(s13, mark, 0.2, y, 0.4, 0.38,
          font_size=14, bold=True, color=mc, align=PP_ALIGN.CENTER)
    txbox(s13, text, 0.65, y, 9.1, 0.38, font_size=13, color=WHITE)

footer(s13)

# ── Слайд 14: РЕКОМЕНДАЦИИ ────────────────────────────────────
make_photo_slide(
    title="ФОТО 5: РЕКОМЕНДАЦИИ / БЕЗОПАСНОСТЬ — примеры конкурентов",
    photos=[
        ("6664644908", "GLITTER G10/G5\n«Наденьте перчатки» крупно\nФото синей перчатки",         "GLITTER"),
        ("9405795820", "GLITTER G5\n«Защитные перчатки»\nМужчина на мойке с пистолетом",           "GLITTER"),
        ("8278220086", "DR.BERG ACE\n«Способ применения» + предупреждения\nТёмный, чёткий текст",  "DR.BERG"),
    ],
    focus_items=[
        ("✓", "Реальное фото человека в перчатках (не иконка) — как у GLITTER"),
        ("✓", "Главное предупреждение КРУПНО: «НЕ НАНОСИТЬ НА ГОРЯЧИЙ КУЗОВ»"),
        ("✓", "«ВЫДЕРЖАТЬ НЕ БОЛЕЕ 2 МИНУТ» — жёлтый акцент"),
        ("✓", "Блок антивозражений: «что будет если...» — закрываем страхи"),
        ("✗", "Не только список ⚠ — нужна визуальная иерархия по важности"),
    ],
    cols=3,
)

# ── Слайд 15: ПРОИЗВОДСТВО / ДОВЕРИЕ ─────────────────────────
make_photo_slide(
    title="ФОТО 6: ПРОИЗВОДСТВО И ДОВЕРИЕ — примеры конкурентов",
    photos=[
        ("9405795888", "GLITTER G5\n«Сделано в России для наших дорог»\nКоллаж лаборатории 2×2",    "GLITTER"),
        ("9405795838", "GLITTER G5\n«51 000+ положительных отзывов»\nКарточки с цитатами и ★",     "GLITTER"),
        ("9405795878", "GLITTER G5\n«Мойка своими руками — как у профи»\nМужчина с канистром",     "GLITTER"),
        ("9405795857", "GLITTER G5\n«Почему мы?» — таблица сравнения\nGLITTER vs обычное средство", "GLITTER"),
    ],
    focus_items=[
        ("✓", "«ПРОИЗВОДИМ САМИ — ОТВЕЧАЕМ ЗА КАЧЕСТВО» — фото цеха ifoam"),
        ("✓", "Флаг РФ + ТУ 20.41.32-003-36603872-2024 — закрываем «подделка?»"),
        ("✓", "Коллаж 2×2: оборудование / лаборатория / розлив / готовый продукт"),
        ("✓", "Если есть отзывы — 2 карточки с цитатами и ★★★★★"),
        ("✗", "НЕ AI-фото лаборатории — только реальные кадры производства ifoam"),
    ],
)

# ── сохранить ─────────────────────────────────────────────────
prs.save(str(pptx_path))
print(f"Готово: {pptx_path}")
print(f"Итого слайдов: {len(prs.slides)}")
shutil.rmtree(tmp_dir, ignore_errors=True)
