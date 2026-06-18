# -*- coding: utf-8 -*-
"""
Генератор ТЗ для дизайнера — ifoam Orange Standard 5кг (Art. 119105)
Структура: аналог примеров All cleaner / Gloss Disk / POLYROL FROST
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Цвета бренда ifoam Orange Standard ──────────────────────────
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)   # почти чёрный фон
ORANGE      = RGBColor(0xFF, 0x6B, 0x00)   # оранжевый — цвет жидкости / бренд-акцент
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY   = RGBColor(0x2D, 0x2D, 0x2D)
ORANGE_LITE = RGBColor(0xFF, 0xA0, 0x40)   # светлый оранжевый для подзаголовков

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(13.33)   # портретный формат 3:4 (как у конкурентов)

BLANK = prs.slide_layouts[6]   # пустой макет


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color: RGBColor):
    """Залить фон слайда цветом."""
    from pptx.util import Emu
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txbox(slide, text, left, top, width, height,
          font_size=24, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, font_name="Arial"):
    """Добавить текстовый блок."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    """Добавить прямоугольник."""
    from pptx.util import Pt as PPt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = PPt(line_width)
    else:
        shape.line.fill.background()
    return shape


def section_label(slide, num, title, top=0.25):
    """Номер + заголовок раздела — оранжевый номер + белый текст."""
    txbox(slide, str(num), 0.3, top, 0.7, 0.7,
          font_size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txbox(slide, title, 1.1, top+0.05, 8.5, 0.65,
          font_size=28, bold=True, color=WHITE)


def bullet_block(slide, items, left, top, width=8.5, font_size=18,
                 bullet="✓", bullet_color=ORANGE, text_color=WHITE, spacing=0.42):
    """Блок с маркерами."""
    for i, item in enumerate(items):
        txbox(slide, bullet, left, top + i * spacing, 0.5, 0.45,
              font_size=font_size, bold=True, color=bullet_color)
        txbox(slide, item, left + 0.5, top + i * spacing, width - 0.5, 0.45,
              font_size=font_size, color=text_color)


def warn_block(slide, items, left, top, font_size=17):
    bullet_block(slide, items, left, top, font_size=font_size,
                 bullet="⚠", bullet_color=RGBColor(0xFF, 0xCC, 0x00), text_color=WHITE)


def divider(slide, top, color=ORANGE):
    """Горизонтальная оранжевая полоска-разделитель."""
    rect(slide, 0.3, top, 9.4, 0.04, fill_color=color)


# ════════════════════════════════════════════════════════════════
# СЛАЙД 1 — ОБЩИЕ ТРЕБОВАНИЯ
# ════════════════════════════════════════════════════════════════
s1 = add_slide()
bg(s1, DARK_GRAY)

txbox(s1, "ОБЩИЕ ТРЕБОВАНИЯ", 0.3, 0.25, 9.4, 0.9,
      font_size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
divider(s1, 1.2)

# Левая колонка — Формат
txbox(s1, "Формат изображений:", 0.3, 1.4, 4.5, 0.5,
      font_size=20, bold=True, color=ORANGE)
fmt_items = [
    "Размер: 1200×1200 px (главное фото)",
    "         1200×1600 px (инфографика / вертикаль)",
    "Формат: JPG / PNG",
    "Качество: не менее 300 dpi",
    "Цветовое пространство: RGB",
]
for i, t in enumerate(fmt_items):
    txbox(s1, t, 0.3, 1.9 + i*0.38, 4.5, 0.4, font_size=16, color=LIGHT_GRAY)

# Правая колонка — Стиль
txbox(s1, "Стиль:", 5.2, 1.4, 4.5, 0.5,
      font_size=20, bold=True, color=ORANGE)
style_items = [
    "Тёмный премиальный фон (чёрный/тёмно-серый)",
    "Акцент — оранжевый (#FF6B00) под цвет жидкости",
    "Крупная белая типографика, читаемая на мобильном",
    "Продукт всегда виден в кадре",
    "Профессиональный, технический образ",
    "Атмосфера: «мойка как у профи, дома»",
]
for i, t in enumerate(style_items):
    txbox(s1, "• " + t, 5.2, 1.9 + i*0.38, 4.5, 0.4, font_size=16, color=LIGHT_GRAY)

divider(s1, 4.25)

# Требования к шрифту
txbox(s1, "Требования к тексту на изображениях:", 0.3, 4.4, 9.4, 0.5,
      font_size=20, bold=True, color=ORANGE)
txt_items = [
    "Заголовки — жирный, без засечек (Arial Bold / Montserrat Bold)",
    "Минимальный размер текста — 24 px (читаемость на мобильном)",
    "Текст не перекрывает логотип и ключевые элементы этикетки",
    "На каждом изображении — не более 2–3 смысловых блоков",
]
for i, t in enumerate(txt_items):
    txbox(s1, "• " + t, 0.3, 4.9 + i*0.38, 9.4, 0.4, font_size=16, color=LIGHT_GRAY)

divider(s1, 6.55)

# Цветовая палитра
txbox(s1, "Цветовая палитра:", 0.3, 6.7, 9.4, 0.5,
      font_size=20, bold=True, color=ORANGE)
palette = [
    ("Основной фон", "#1A1A1A / #2D2D2D", BLACK),
    ("Акцент (бренд)", "#FF6B00", ORANGE),
    ("Текст", "#FFFFFF", WHITE),
    ("Вспомогательный", "#FFA040", ORANGE_LITE),
]
for i, (name, hex_val, col) in enumerate(palette):
    x = 0.3 + i * 2.35
    rect(s1, x, 7.25, 2.0, 0.45, fill_color=col,
         line_color=RGBColor(0x80,0x80,0x80), line_width=1)
    txbox(s1, f"{name}\n{hex_val}", x, 7.75, 2.0, 0.65, font_size=13, color=LIGHT_GRAY)

divider(s1, 8.55)

# Запрещено
txbox(s1, "Запрещено:", 0.3, 8.7, 9.4, 0.5,
      font_size=20, bold=True, color=RGBColor(0xFF,0x44,0x44))
forbidden = [
    "Пересечение текста с основным изображением товара",
    "Использование зелёного цвета (ассоциация с конкурентами Grass)",
    "Светлый / белый фон — продукт теряется на полке",
    "Мелкий нечитаемый текст поверх сложного фона",
]
for i, t in enumerate(forbidden):
    txbox(s1, "✗  " + t, 0.3, 9.2 + i*0.38, 9.4, 0.4,
          font_size=16, color=RGBColor(0xFF,0x88,0x88))

txbox(s1, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# СЛАЙД 2 — ГЛАВНОЕ ФОТО
# ════════════════════════════════════════════════════════════════
s2 = add_slide()
bg(s2, BLACK)

section_label(s2, 1, "ГЛАВНОЕ ФОТО", top=0.25)
divider(s2, 1.05)

txbox(s2, "Фон:", 0.3, 1.2, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
bg_items = [
    "Основной цвет: насыщенный чёрный (#1A1A1A)",
    "Дополнение: лёгкий оранжевый градиент в нижней трети — перекличка с цветом жидкости",
    "Ассоциация: профессиональная автохимия, сила, доверие",
]
for i, t in enumerate(bg_items):
    txbox(s2, "• " + t, 0.3, 1.65 + i*0.38, 9.4, 0.4, font_size=16, color=LIGHT_GRAY)

txbox(s2, "Композиция:", 0.3, 2.9, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
comp_items = [
    "Канистра 5 кг по центру, лёгкий наклон 10–15° — динамика",
    "Этикетка Orange Standard чётко читается, повёрнута к зрителю",
    "Крышка закрыта, горловина запаяна — акцент на защите от подделок",
    "Лёгкий блик/свет сверху справа — объём, «живость» упаковки",
    "Фон без лишних предметов — фокус только на продукте",
]
for i, t in enumerate(comp_items):
    txbox(s2, "• " + t, 0.3, 3.35 + i*0.38, 9.4, 0.4, font_size=16, color=LIGHT_GRAY)

txbox(s2, "Текстовые блоки на фото:", 0.3, 5.4, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)

# Текстовые блоки — как у конкурентов
text_blocks = [
    ("АВТОШАМПУНЬ", "Основной заголовок — крупно сверху, белый жирный"),
    ("БЕСКОНТАКТНАЯ МОЙКА", "Подзаголовок — оранжевый, чуть меньше"),
    ("pH 12  |  концентрат 1:30–1:40", "Параметры — белый, строго под подзаголовком"),
    ("5 кг", "Объём — круглый бейдж нижний левый угол, оранжевый фон"),
    ("ОТ ПРОИЗВОДИТЕЛЯ", "Доп. метка нижний правый угол, небольшой шрифт"),
]
for i, (label, desc) in enumerate(text_blocks):
    txbox(s2, f'"{label}"', 0.3, 5.9 + i*0.52, 3.5, 0.45,
          font_size=16, bold=True, color=ORANGE)
    txbox(s2, desc, 3.9, 5.9 + i*0.52, 5.8, 0.45,
          font_size=16, color=LIGHT_GRAY)

divider(s2, 8.65)
txbox(s2, "Справка для дизайнера — что делают лидеры (GLITTER G10, GRASS Balance):",
      0.3, 8.8, 9.4, 0.45, font_size=15, bold=True, color=ORANGE_LITE)
notes = [
    "GLITTER G10: чистый продакт-шот на светлом, весь текст на этикетке, минимум надписей на фото",
    "GRASS Balance: яркий жёлтый канистр + «Бренд №1» бейдж — доверие через статус",
    "ifoam: должен взять лучшее — тёмный фон (как DR.BERG/GLITTER) + яркий оранжевый акцент",
]
for i, t in enumerate(notes):
    txbox(s2, "→ " + t, 0.3, 9.3 + i*0.42, 9.4, 0.4, font_size=14,
          color=RGBColor(0xCC,0xCC,0xCC), italic=True)

txbox(s2, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# СЛАЙД 3 — ХАРАКТЕРИСТИКИ / УТП
# ════════════════════════════════════════════════════════════════
s3 = add_slide()
bg(s3, BLACK)

section_label(s3, 2, 'ФОТО "ХАРАКТЕРИСТИКИ / УТП"', top=0.25)
divider(s3, 1.05)

txbox(s3, "Фон:", 0.3, 1.2, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
txbox(s3, "Тёмный с оранжевыми пузырями пены или брызгами воды — атмосфера мойки. "
     "Канистр в правом нижнем углу, занимает 30–35% кадра.",
      0.3, 1.62, 9.4, 0.6, font_size=16, color=LIGHT_GRAY)

txbox(s3, "Текстовые блоки (✓ маркеры, белый + оранжевый):", 0.3, 2.35, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)

bullet_items = [
    "Щелочная формула pH 12 — растворяет самые стойкие загрязнения",
    "Концентрат 1:30–1:40 — 5 кг = 150–200 литров рабочего раствора",
    "Не повреждает ЛКП, хром, пластик и резину",
    "Работает с мягкой и жёсткой водой",
    "Густая активная пена держится 1–2 минуты",
    "Легко смывается водой без разводов",
]
bullet_block(s3, bullet_items, 0.3, 2.88, width=9.2, font_size=17, spacing=0.44)

divider(s3, 5.65)

txbox(s3, "Акцентный блок (выделить рамкой / плашкой):", 0.3, 5.82, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
rect(s3, 0.3, 6.32, 9.4, 1.6, fill_color=RGBColor(0x2A,0x14,0x00),
     line_color=ORANGE, line_width=2)
txbox(s3, "pH 12", 1.0, 6.45, 2.2, 0.75,
      font_size=40, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txbox(s3, "Щелочная\nформула", 1.0, 7.15, 2.2, 0.6,
      font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txbox(s3, "1:30–1:40", 3.8, 6.45, 2.2, 0.75,
      font_size=40, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txbox(s3, "Пеногенератор", 3.8, 7.15, 2.2, 0.6,
      font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txbox(s3, "1:3–1:4", 6.6, 6.45, 2.2, 0.75,
      font_size=40, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txbox(s3, "Пенокомплект", 6.6, 7.15, 2.2, 0.6,
      font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

divider(s3, 8.12)

txbox(s3, "Справка — что используют конкуренты на этом слайде:", 0.3, 8.28, 9.4, 0.45,
      font_size=15, bold=True, color=ORANGE_LITE)
notes3 = [
    "GLITTER G10: «работает с мягкой и жёсткой водой / густая пена / не оставляет разводов» + иконки",
    "DR.BERG ACE: «активная пена / эффективное удаление / плотная пена / безопасно» на тёмном фоне",
    "ifoam: добавить уникальное — pH 12 и пропорции разбавления крупно (это наше реальное УТП)",
]
for i, t in enumerate(notes3):
    txbox(s3, "→ " + t, 0.3, 8.78 + i*0.42, 9.4, 0.4, font_size=14,
          color=RGBColor(0xCC,0xCC,0xCC), italic=True)

txbox(s3, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# СЛАЙД 4 — ПРИМЕНЕНИЕ (инструкция)
# ════════════════════════════════════════════════════════════════
s4 = add_slide()
bg(s4, BLACK)

section_label(s4, 3, 'ФОТО "ПРИМЕНЕНИЕ"', top=0.25)
divider(s4, 1.05)

txbox(s4, "Формат: пошаговая инструкция с реальными фотографиями или иконками (4 шага)",
      0.3, 1.2, 9.4, 0.55, font_size=17, color=LIGHT_GRAY)

txbox(s4, "Заголовок слайда (крупно, сверху):",
      0.3, 1.85, 9.4, 0.45, font_size=19, bold=True, color=ORANGE)
txbox(s4, '"4 ШАГА ДО ЧИСТОГО КУЗОВА"  или  "ИНСТРУКЦИЯ ПО ПРИМЕНЕНИЮ"',
      0.3, 2.3, 9.4, 0.45, font_size=17, bold=True, color=WHITE)

divider(s4, 2.9)

steps = [
    ("1", "ОСТУДИТЕ КУЗОВ",
     "Сбейте сильные загрязнения струёй воды\nпри необходимости"),
    ("2", "РАЗВЕДИТЕ СОСТАВ",
     "Пеногенератор: 20–25 г/л (1:30–1:40)\nПенокомплект: 200–250 г/л (1:3–1:4)"),
    ("3", "НАНЕСИТЕ ПЕНУ",
     "Равномерно покройте весь кузов снизу вверх\nВыдержите 1–2 минуты"),
    ("4", "СМОЙТЕ ВОДОЙ",
     "Сбейте пену под давлением снизу вверх\nНЕ допускайте высыхания на поверхности"),
]

for i, (num, title, desc) in enumerate(steps):
    top = 3.05 + i * 2.3
    rect(s4, 0.3, top, 9.4, 2.1,
         fill_color=RGBColor(0x22,0x22,0x22),
         line_color=ORANGE if i == 2 else RGBColor(0x44,0x44,0x44),
         line_width=1.5)
    txbox(s4, num, 0.45, top + 0.15, 0.85, 1.1,
          font_size=52, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txbox(s4, title, 1.45, top + 0.2, 7.9, 0.6,
          font_size=22, bold=True, color=WHITE)
    txbox(s4, desc, 1.45, top + 0.82, 7.9, 0.95,
          font_size=16, color=LIGHT_GRAY)

txbox(s4, "⚠  НЕ ДОПУСКАЙТЕ ВЫСЫХАНИЯ НА ПОВЕРХНОСТИ",
      0.3, 12.2, 9.4, 0.5,
      font_size=17, bold=True, color=RGBColor(0xFF,0xCC,0x00),
      align=PP_ALIGN.CENTER)

txbox(s4, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# СЛАЙД 5 — ЭКОНОМИКА
# ════════════════════════════════════════════════════════════════
s5 = add_slide()
bg(s5, BLACK)

section_label(s5, 4, 'ФОТО "ЭКОНОМИКА"', top=0.25)
divider(s5, 1.05)

txbox(s5, "Концепция: финансовая выгода — уникальный слайд, которого НЕТ у конкурентов",
      0.3, 1.2, 9.4, 0.55, font_size=16, italic=True, color=ORANGE_LITE)

txbox(s5, "Заголовок слайда (крупно, сверху):",
      0.3, 1.85, 9.4, 0.45, font_size=19, bold=True, color=ORANGE)
txbox(s5, '"5 КГ = 20 АВТОМОБИЛЕЙ"  или  "СТОИМОСТЬ 1 МОЙКИ — ОТ 12 РУБЛЕЙ"',
      0.3, 2.3, 9.4, 0.55, font_size=17, bold=True, color=WHITE)

divider(s5, 3.0)

# Три цифровых блока
stats = [
    ("5 кг", "концентрата"),
    ("150–200 л", "рабочего раствора"),
    ("~12 руб.", "стоимость 1 мойки"),
]
for i, (val, label) in enumerate(stats):
    x = 0.3 + i * 3.1
    rect(s5, x, 3.15, 2.9, 2.0, fill_color=RGBColor(0x2A,0x14,0x00),
         line_color=ORANGE, line_width=2)
    txbox(s5, val, x, 3.25, 2.9, 0.9,
          font_size=30, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txbox(s5, label, x, 4.1, 2.9, 0.6,
          font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txbox(s5, "→", 0.3, 4.0, 0.5, 0.6, font_size=30, bold=True, color=ORANGE)

divider(s5, 5.35)

txbox(s5, "Сравнение (визуальная таблица):", 0.3, 5.5, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)

# Таблица: заголовок
rect(s5, 0.3, 6.0, 9.4, 0.55, fill_color=ORANGE)
txbox(s5, "Способ мойки", 0.35, 6.05, 3.0, 0.45, font_size=15, bold=True, color=BLACK)
txbox(s5, "Стоимость", 3.4, 6.05, 3.0, 0.45, font_size=15, bold=True, color=BLACK)
txbox(s5, "ifoam 5 кг", 6.5, 6.05, 3.0, 0.45, font_size=15, bold=True, color=BLACK)

table_rows = [
    ("Автомойка (мойщик)", "300–600 руб./авто", "—"),
    ("Мойка самообслуживания", "100–200 руб./авто", "—"),
    ("ifoam Orange Standard 5 кг", "~12 руб./авто", "✓"),
]
row_colors = [RGBColor(0x28,0x28,0x28), RGBColor(0x22,0x22,0x22), RGBColor(0x33,0x18,0x00)]
for j, (c1, c2, c3) in enumerate(table_rows):
    top = 6.6 + j * 0.6
    rect(s5, 0.3, top, 9.4, 0.55, fill_color=row_colors[j])
    txbox(s5, c1, 0.35, top+0.05, 3.0, 0.45, font_size=15, color=WHITE)
    txbox(s5, c2, 3.4, top+0.05, 3.0, 0.45, font_size=15, color=LIGHT_GRAY)
    col3 = ORANGE if c3 == "✓" else LIGHT_GRAY
    txbox(s5, c3, 6.5, top+0.05, 3.0, 0.45, font_size=15, bold=(c3=="✓"), color=col3)

txbox(s5, "Экономия до 97% по сравнению с автомойкой", 0.3, 8.45, 9.4, 0.55,
      font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

divider(s5, 9.15)
txbox(s5, "Дополнительный текстовый блок (подтверждение):", 0.3, 9.3, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
bullet_block(s5, [
    "На одно авто расход: 30–50 г концентрата",
    "1 упаковка = минимум 20 полных моек легкового автомобиля",
    "Выгоднее фасовки 1 кг при том же качестве",
], 0.3, 9.82, width=9.2, font_size=16, spacing=0.45)

txbox(s5, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# СЛАЙД 6 — РЕКОМЕНДАЦИИ / БЕЗОПАСНОСТЬ
# ════════════════════════════════════════════════════════════════
s6 = add_slide()
bg(s6, BLACK)

section_label(s6, 5, 'КАРТОЧКА "РЕКОМЕНДАЦИИ"', top=0.25)
divider(s6, 1.05)

txbox(s6, "Фон: тёмный с размытым фото автомобиля или гаражного пространства",
      0.3, 1.2, 9.4, 0.45, font_size=16, color=LIGHT_GRAY, italic=True)

txbox(s6, "Меры предосторожности:", 0.3, 1.75, 9.4, 0.45,
      font_size=20, bold=True, color=ORANGE)
warn_block(s6, [
    "Не наносить на нагретые поверхности (температура кузова выше 40°С)",
    "Не допускать высыхания на поверхности — смывать не позднее 2 минут",
    "Избегать попадания в глаза и на кожу — использовать перчатки",
    "Не смешивать с кислотными средствами",
    "Беречь от детей и животных",
    "Хранить при температуре от 0°С до +30°С",
], 0.3, 2.3)

divider(s6, 5.6)

txbox(s6, "Блок «Что будет при нарушении» (антивозражение):", 0.3, 5.75, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
anti = [
    ("Нанесли на горячий кузов?",
     "Средство высохнет раньше смывания — оставит белесые следы. Решение: работайте на холодном кузове."),
    ("Держали дольше 2 минут?",
     "Возможны разводы. Решение: контролируйте время, смывайте своевременно."),
    ("Не разбавили концентрат?",
     "Избыточная щёлочь не улучшает результат — только расход. Соблюдайте пропорцию 1:30–1:40."),
]
for i, (q, a) in enumerate(anti):
    top = 6.28 + i * 1.8
    rect(s6, 0.3, top, 9.4, 1.65,
         fill_color=RGBColor(0x22,0x22,0x22),
         line_color=RGBColor(0xFF,0xCC,0x00), line_width=1)
    txbox(s6, "❓  " + q, 0.45, top + 0.1, 9.1, 0.5,
          font_size=16, bold=True, color=RGBColor(0xFF,0xCC,0x00))
    txbox(s6, a, 0.45, top + 0.65, 9.1, 0.85,
          font_size=15, color=LIGHT_GRAY)

txbox(s6, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# СЛАЙД 7 — ПРОИЗВОДСТВО / ДОВЕРИЕ
# ════════════════════════════════════════════════════════════════
s7 = add_slide()
bg(s7, BLACK)

section_label(s7, 6, 'КАРТОЧКА "ПРОИЗВОДСТВО И ДОВЕРИЕ"', top=0.25)
divider(s7, 1.05)

txbox(s7, "Концепция: «Сделано в России — мы отвечаем за качество»  (как у GLITTER G5)",
      0.3, 1.2, 9.4, 0.45, font_size=16, italic=True, color=ORANGE_LITE)

txbox(s7, "Заголовок слайда (крупно):", 0.3, 1.75, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
txbox(s7, '"ПРОИЗВОДИМ САМИ — ОТВЕЧАЕМ ЗА КАЧЕСТВО"',
      0.3, 2.2, 9.4, 0.55, font_size=19, bold=True, color=WHITE)
txbox(s7, "Подзаголовок: «ifoam AUTO — российский производитель автохимии»",
      0.3, 2.78, 9.4, 0.45, font_size=16, color=LIGHT_GRAY)

divider(s7, 3.35)

txbox(s7, "Фото-коллаж (4 фрагмента в сетке 2×2):", 0.3, 3.5, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)

photo_labels = [
    ("Производственный цех /\nоборудование ifoam", "Верхний левый"),
    ("Лаборант / контроль\nкачества продукта", "Верхний правый"),
    ("Канистра Orange Standard\nна линии розлива", "Нижний левый"),
    ("Стенд с продукцией\nлинейки ifoam AUTO", "Нижний правый"),
]
positions = [(0.3, 4.05), (5.05, 4.05), (0.3, 7.05), (5.05, 7.05)]
for (px, py), (label, pos) in zip(positions, photo_labels):
    rect(s7, px, py, 4.6, 2.7,
         fill_color=RGBColor(0x28,0x28,0x28),
         line_color=ORANGE, line_width=1.5)
    txbox(s7, f"[ФОТО]\n{label}", px + 0.15, py + 0.7, 4.3, 1.1,
          font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    txbox(s7, pos, px + 0.15, py + 0.1, 4.3, 0.4,
          font_size=12, color=ORANGE, align=PP_ALIGN.CENTER)

divider(s7, 9.9)

txbox(s7, "Нижний баннер (поверх фото / под сеткой):", 0.3, 10.05, 9.4, 0.45,
      font_size=19, bold=True, color=ORANGE)
rect(s7, 0.3, 10.55, 9.4, 1.0,
     fill_color=RGBColor(0x2A,0x14,0x00),
     line_color=ORANGE, line_width=2)
txbox(s7, "🇷🇺  ifoam AUTO  |  ТУ 20.41.32-003-36603872-2024  |  Сделано в России",
      0.3, 10.7, 9.4, 0.6, font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s7, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# СЛАЙД 8 — ТЕХНИЧЕСКИЕ ТРЕБОВАНИЯ + ВАЖНЫЕ АКЦЕНТЫ
# ════════════════════════════════════════════════════════════════
s8 = add_slide()
bg(s8, DARK_GRAY)

txbox(s8, "ТЕХНИЧЕСКИЕ ТРЕБОВАНИЯ", 0.3, 0.25, 9.4, 0.8,
      font_size=34, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
divider(s8, 1.1)

# Шрифты
txbox(s8, "Шрифты:", 0.3, 1.28, 9.4, 0.45,
      font_size=20, bold=True, color=ORANGE)
font_items = [
    "Заголовки: жирный, без засечек — Arial Bold / Montserrat Bold / Impact",
    "Основной текст: чёткий, контрастный — Arial / Roboto",
    "Минимальный размер: 24 px (читаемость на мобильном 375 px)",
    "Не использовать декоративные и курсивные шрифты для ключевых надписей",
]
for i, t in enumerate(font_items):
    txbox(s8, "• " + t, 0.3, 1.78 + i*0.38, 9.4, 0.4, font_size=16, color=LIGHT_GRAY)

divider(s8, 3.35)

# Цвета
txbox(s8, "Цветовая палитра:", 0.3, 3.52, 9.4, 0.45,
      font_size=20, bold=True, color=ORANGE)
colors_table = [
    ("Основной фон", "#1A1A1A", "Все слайды"),
    ("Акцент бренда", "#FF6B00", "Заголовки, маркеры, рамки"),
    ("Текст основной", "#FFFFFF", "Все подписи"),
    ("Текст вторичный", "#F2F2F2", "Описания, пояснения"),
    ("Предупреждения", "#FFCC00", "⚠ блоки"),
    ("Запрещено", "#FF4444", "✗ блоки"),
]
for i, (name, hex_v, use) in enumerate(colors_table):
    txbox(s8, name, 0.3, 4.02 + i*0.38, 2.8, 0.38, font_size=15, color=LIGHT_GRAY)
    txbox(s8, hex_v, 3.2, 4.02 + i*0.38, 2.0, 0.38, font_size=15, bold=True, color=ORANGE)
    txbox(s8, use,   5.3, 4.02 + i*0.38, 4.4, 0.38, font_size=15, color=LIGHT_GRAY)

divider(s8, 6.35)

# Важные акценты
txbox(s8, "ВАЖНЫЕ АКЦЕНТЫ — УТП ifoam Orange Standard:", 0.3, 6.52, 9.4, 0.55,
      font_size=20, bold=True, color=ORANGE)
accents = [
    "pH 12 — ключевой технический параметр, выносить на каждую инфографику",
    "Цвет жидкости оранжевый — использовать как фирменный визуальный акцент",
    "Конкуренты используют синий (GLITTER) и жёлтый (GRASS) — нам нужен чёткий контраст",
    "Экономика (12 руб./мойка) — уникальный слайд, аналогов у конкурентов нет",
    "«От производителя» и российский флаг — закрываем возражение о подделках",
    "Не копировать зелёный цвет — GRASS занял эту нишу, путаница недопустима",
]
bullet_block(s8, accents, 0.3, 7.15, font_size=15, spacing=0.42)

divider(s8, 9.75)
txbox(s8, "СОСТАВ ПАКЕТА ФОТОГРАФИЙ (итого 6 изображений):", 0.3, 9.92, 9.4, 0.45,
      font_size=18, bold=True, color=ORANGE)
pack = [
    "1. Главное фото (продакт-шот) — 1200×1200 px",
    "2. Характеристики / УТП — 1200×1600 px",
    "3. Инструкция по применению — 1200×1600 px",
    "4. Экономика (5 кг = 20 авто) — 1200×1600 px",
    "5. Рекомендации и безопасность — 1200×1600 px",
    "6. Производство и доверие — 1200×1600 px",
]
for i, t in enumerate(pack):
    txbox(s8, t, 0.3, 10.42 + i*0.38, 9.4, 0.38, font_size=15, color=LIGHT_GRAY)

txbox(s8, "ifoam AUTO  |  Orange Standard 5кг  |  Art. 119105",
      0.3, 12.9, 9.4, 0.35, font_size=12,
      color=RGBColor(0x80,0x80,0x80), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# СОХРАНЕНИЕ
# ════════════════════════════════════════════════════════════════
import sys
sys.path.insert(0, "C:/Users/1/Downloads/ozon_prompts_v2")
from save_helpers import get_product_folder
folder = get_product_folder("ifoam", "OrangeStandard5kg", "119105")
out = folder / "07_designer_tz_119105_v2.pptx"
prs.save(str(out))
print(f"Saved: {out}")
print(f"Слайдов: {len(prs.slides)}")
